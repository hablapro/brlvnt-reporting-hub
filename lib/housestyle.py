#!/usr/bin/env python3
"""Berelvant FOREX.com house style for performance-review decks.

Single source of the design tokens and slide primitives. Every monthly deck
builder imports from here; nothing re-declares a color, a font, or a footer.

Why this exists: the June 2026 cycle ran two ~700-line builders (GGMI, GCG)
that were ~50% identical copies, built in an old palette, then converted to
the house palette by a separate restyle pass. Three places to keep in sync,
so style drifted and the footer page-count bug shipped twice. This module is
the deduplicated version, in house tokens from the start. No restyle step.

The spec these tokens implement is docs/DESIGN-SYSTEM.md. That document's
coordinates describe a 10 x 5.625in grid; this module uses the 13.333 x 7.5in
grid the shipped Python builders use. Multiply spec coordinates by 1.3333.

Usage:

    from housestyle import Deck, CORAL, GREEN, GOLD

    d = Deck(entity='GGMI (LATAM)', month='July 2026', n_slides=15)
    s = d.cover(kicker='QUARTERLY BUSINESS REVIEW',
                title='GGMI (LATAM) - Paid Media',
                subtitle='July 2026 Performance',
                channel_line='Bing Ads . Meta . Azerion . Quantcast',
                footer='Reporting period: July 1-31, 2026 | Currency: USD')
    s = d.content('EXECUTIVE SUMMARY', 'The headline for this slide.')
    d.table(s, 0.5, 1.9, 7.4, 2.0, rows, col_widths=[...], total_last=True)
    d.card(s, 8.1, 1.9, 4.7, 2.4, 'WHAT THE DATA SAYS', [...], GREEN)
    d.save('/path/to/deck.pptx')

Run `python3 lib/housestyle.py` to build a smoke-test deck exercising every
primitive, then render it to check the style before trusting a real build.
"""
import os

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import (XL_CHART_TYPE, XL_LABEL_POSITION,
                             XL_LEGEND_POSITION, XL_MARKER_STYLE)
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---------------------------------------------------------------- tokens ----
# docs/DESIGN-SYSTEM.md is the human-readable spec for these. Do not add
# colors here without updating that file; one navy dominates, coral is the
# single sharp accent, gold and green carry read/recommend semantics only.
DEEP = RGBColor(0x0F, 0x15, 0x35)    # cover bg, closing strip, "measurement"
NAVY = RGBColor(0x1E, 0x27, 0x61)    # title bands, table headers, values
CORAL = RGBColor(0xF9, 0x61, 0x67)   # band rule, card left edge, "critical"
GOLD = RGBColor(0xE9, 0xB4, 0x4C)    # "what we recommend", "high"
GREEN = RGBColor(0x2E, 0x7D, 0x57)   # "what the data shows", positive deltas
ICE = RGBColor(0xCA, 0xDC, 0xFC)     # cover subtitle, table Total row fill
MUTED = RGBColor(0x5A, 0x60, 0x72)   # stat labels, footer, footnotes
# MUTED on DEEP is ~2.3:1 and unreadable on the cover. This is the same muted
# role lightened for dark fills. Verified in a real PowerPoint render.
MUTED_DARK = RGBColor(0x9A, 0xA4, 0xBA)
BORDER = RGBColor(0xD0, 0xD4, 0xDC)  # card and table borders
LIGHT = RGBColor(0xF9, 0xFA, 0xFC)   # alternating rows, blocker card body
INK = RGBColor(0x2B, 0x31, 0x47)     # body text
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HEAD_FONT = 'Georgia'   # headers, titles, section labels, read-box titles
BODY_FONT = 'Calibri'   # body, bullets, tables, stat labels/values, footer

W, H = 13.333, 7.5      # 16:9
MARGIN = 0.5
FOOTER_Y = 7.2          # keep content above ~7.05 so it clears the footer

_ASSETS = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       '..', 'assets')
FOREX_LOGO = os.path.join(_ASSETS, 'forex_logo.png')
BERELVANT_LOGO = os.path.join(_ASSETS, 'berelvant_logo.png')


class Deck:
    """A house-style deck. One instance per deliverable.

    n_slides is the denominator in the "n / total" footer. Set it to the
    final slide count up front; call verify() before save() and it will
    raise if the real count drifted, which is the check that catches the
    footer-numbering bug that shipped twice in the June cycle.
    """

    def __init__(self, entity, month, n_slides, footer_label=None):
        self.entity = entity
        self.month = month
        self.n_slides = n_slides
        self.footer_label = footer_label or (
            f'{entity}   |   {month} Performance Review   |   Berelvant')
        self.prs = Presentation()
        self.prs.slide_width = Inches(W)
        self.prs.slide_height = Inches(H)
        self._blank = self.prs.slide_layouts[6]
        self._page = 0

    # ------------------------------------------------------------ shapes ----
    def rect(self, slide, x, y, w, h, fill, border=None):
        sh = slide.shapes.add_shape(1, Inches(x), Inches(y),
                                    Inches(w), Inches(h))
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
        if border is None:
            sh.line.fill.background()
        else:
            sh.line.color.rgb = border
            sh.line.width = Pt(0.75)
        sh.shadow.inherit = False
        return sh

    def text(self, slide, x, y, w, h, runs, size=9.5, bold=False, color=INK,
             font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             wrap=True, space=6):
        """runs: a string, or a list of strings rendered as paragraphs."""
        tb = slide.shapes.add_textbox(Inches(x), Inches(y),
                                     Inches(w), Inches(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.vertical_anchor = anchor
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        paras = [runs] if isinstance(runs, str) else runs
        for i, ptext in enumerate(paras):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align
            if i > 0:
                p.space_before = Pt(space)
            r = p.add_run()
            r.text = ptext
            r.font.name = font
            r.font.size = Pt(size)
            r.font.bold = bold
            r.font.color.rgb = color
        return tb

    # ------------------------------------------------------------ slides ----
    def cover(self, kicker, title, subtitle, channel_line, footer):
        """Cover slide. Deep navy, FOREX logo top-left, no page footer."""
        s = self.prs.slides.add_slide(self._blank)
        self._page += 1
        self.rect(s, 0, 0, W, H, DEEP)
        self.rect(s, MARGIN, 2.2, 0.1, 2.5, CORAL)
        self._picture(s, FOREX_LOGO, MARGIN, 0.5, 2.5, 0.71)
        self._picture(s, BERELVANT_LOGO, W - 2.0, H - 1.0, 1.4, 0.4)
        self.text(s, 0.9, 2.2, 11.0, 0.4, kicker, 14, True, CORAL, HEAD_FONT)
        self.text(s, 0.9, 2.7, 12.0, 1.0, title, 36, True, WHITE, HEAD_FONT)
        self.text(s, 0.9, 3.8, 12.0, 0.6, subtitle, 22, False, ICE, HEAD_FONT)
        self.text(s, 0.9, 4.6, 12.0, 0.4, channel_line, 12, False, MUTED_DARK)
        self.text(s, 0.9, 6.7, 12.0, 0.4, footer, 10, False, MUTED_DARK)
        return s

    def content(self, section, headline):
        """Content slide: navy title band + coral rule + headline + footer."""
        s = self.prs.slides.add_slide(self._blank)
        self._page += 1
        self.rect(s, 0, 0, W, 0.69, NAVY)          # 0.92in on the 10in grid
        self.rect(s, 0, 0.69, W, 0.05, CORAL)
        self.text(s, MARGIN, 0.18, 9.0, 0.35, section, 11, True, WHITE,
                  HEAD_FONT)
        self._picture(s, BERELVANT_LOGO, W - 1.6, 0.16, 1.05, 0.3)
        self.text(s, MARGIN, 0.95, 12.3, 0.75, headline, 20, True, NAVY,
                  HEAD_FONT)
        self._footer(s)
        return s

    def _footer(self, s):
        self.rect(s, 0, FOOTER_Y, W, 0.012, BORDER)
        self.text(s, MARGIN, FOOTER_Y + 0.08, 8.0, 0.22, self.footer_label,
                  8, False, MUTED)
        self.text(s, W - MARGIN - 1.0, FOOTER_Y + 0.08, 1.0, 0.22,
                  f'{self._page} / {self.n_slides}', 8, False, MUTED,
                  align=PP_ALIGN.RIGHT)

    def _picture(self, s, path, x, y, w, h):
        if os.path.exists(path):
            s.shapes.add_picture(path, Inches(x), Inches(y),
                                 Inches(w), Inches(h))

    # -------------------------------------------------------- components ----
    def tile(self, s, x, y, label, value, sub=None, w=1.9):
        """Stat card. Four per row. Coral left edge, navy value."""
        h = 1.1 if sub else 1.0
        self.rect(s, x, y, w, h, WHITE, border=BORDER)
        self.rect(s, x, y, 0.06, h, CORAL)
        self.text(s, x + 0.2, y + 0.05, w - 0.3, 0.3, label.upper(), 8, True,
                  MUTED)
        self.text(s, x + 0.2, y + 0.32, w - 0.3, 0.5, value, 15, True, NAVY)
        if sub:
            self.text(s, x + 0.2, y + 0.82, w - 0.3, 0.24, sub, 8, False,
                      MUTED)

    def card(self, s, x, y, w, h, title, body_paras, accent=CORAL):
        """Read / recommend card. Accent sets the semantics:
        GREEN = what the data shows, GOLD = what we recommend,
        CORAL = read / priority / critical, DEEP = measurement to close."""
        self.rect(s, x, y, w, h, WHITE, border=BORDER)
        self.rect(s, x, y, 0.07, h, accent)
        self.text(s, x + 0.24, y + 0.18, w - 0.44, 0.3, title.upper(), 10,
                  True, NAVY, HEAD_FONT)
        self.text(s, x + 0.24, y + 0.55, w - 0.44, h - 0.75, body_paras, 9.5,
                  False, INK)

    def blocker(self, s, x, y, w, h, label, body_paras, color=CORAL):
        """Blocker card: light body under a full-width colored header band.

        Label ink is contrast-picked, not fixed white: white on gold is
        ~2.1:1 and illegible on a projector. Dark navy on gold, white on
        coral / green / deep.
        """
        self.rect(s, x, y, w, h, LIGHT, border=BORDER)
        self.rect(s, x, y, w, 0.32, color)
        ink = DEEP if color in (GOLD, ICE) else WHITE
        self.text(s, x + 0.16, y + 0.06, w - 0.32, 0.24, label.upper(), 9.5,
                  True, ink, HEAD_FONT)
        self.text(s, x + 0.16, y + 0.44, w - 0.32, h - 0.6, body_paras, 9,
                  False, INK)

    def strip(self, s, x, y, w, h, takeaway):
        """Closing takeaway strip: deep navy, one white Georgia line."""
        self.rect(s, x, y, w, h, DEEP)
        self.text(s, x + 0.3, y, w - 0.6, h, takeaway, 12, True, WHITE,
                  HEAD_FONT, anchor=MSO_ANCHOR.MIDDLE)

    def note(self, s, x, y, w, h, footnote):
        """Small muted footnote, e.g. the no-blended-total disclosure."""
        return self.text(s, x, y, w, h, footnote, 7.5, False, MUTED)

    def table(self, s, x, y, w, h, rows, col_widths=None, total_last=False):
        """Navy header row, alternating body rows, ice-blue Total row.

        A cell is a string, or a dict {t, c, b, a} to override color, bold
        and alignment (used for positive MoM deltas in green). Put a literal
        dash in any Total cell that must not be summed - conversions and CPA
        are never blended across channels (docs/DOCTRINE.md).
        """
        n_r, n_c = len(rows), len(rows[0])
        gfx = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
        tbl = gfx.table
        tbl.first_row = False
        tbl.horz_banding = False
        if col_widths:
            for ci, cw in enumerate(col_widths):
                tbl.columns[ci].width = Inches(cw)
        for ri, row in enumerate(rows):
            is_total = total_last and ri == n_r - 1
            for ci, val in enumerate(row):
                spec = val if isinstance(val, dict) else {'t': val}
                cell = tbl.cell(ri, ci)
                cell.margin_left = cell.margin_right = Inches(0.06)
                cell.margin_top = cell.margin_bottom = Inches(0.03)
                cell.fill.solid()
                if ri == 0:
                    cell.fill.fore_color.rgb = NAVY
                elif is_total:
                    cell.fill.fore_color.rgb = ICE
                else:
                    cell.fill.fore_color.rgb = WHITE if ri % 2 else LIGHT
                p = cell.text_frame.paragraphs[0]
                r = p.add_run()
                r.text = str(spec['t'])
                r.font.name = BODY_FONT
                r.font.size = Pt(8.5)
                r.font.bold = spec.get(
                    'b', (ri == 0) or (ci == 0) or is_total)
                r.font.color.rgb = spec.get(
                    'c', WHITE if ri == 0 else NAVY if is_total else INK)
                p.alignment = spec.get(
                    'a', PP_ALIGN.LEFT if ci == 0 else PP_ALIGN.RIGHT)
        return gfx

    def line_chart(self, s, x, y, w, h, categories, series, colors=None,
                   label_series=None):
        """series: list of (name, values). label_series: index to direct-label."""
        colors = colors or [NAVY, CORAL, GREEN, GOLD]
        cd = CategoryChartData()
        cd.categories = categories
        for name, vals in series:
            cd.add_series(name, vals)
        gfx = s.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x),
                                 Inches(y), Inches(w), Inches(h), cd)
        ch = gfx.chart
        ch.has_title = False
        ch.font.size = Pt(8)
        ch.font.name = BODY_FONT
        ch.font.color.rgb = INK
        ch.has_legend = len(series) > 1
        if ch.has_legend:
            ch.legend.position = XL_LEGEND_POSITION.BOTTOM
            ch.legend.include_in_layout = False
            ch.legend.font.size = Pt(8)
        va = ch.value_axis
        va.has_major_gridlines = True
        va.major_gridlines.format.line.color.rgb = BORDER
        va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.number_format = '#,##0'
        va.tick_labels.number_format_is_linked = False
        va.format.line.fill.background()
        ca = ch.category_axis
        ca.has_major_gridlines = False
        ca.format.line.color.rgb = BORDER
        for i, ser in enumerate(ch.series):
            ser.smooth = False
            ser.format.line.color.rgb = colors[i % len(colors)]
            ser.format.line.width = Pt(2.25)
            ser.marker.style = XL_MARKER_STYLE.CIRCLE
            ser.marker.size = 6
            ser.marker.format.fill.solid()
            ser.marker.format.fill.fore_color.rgb = colors[i % len(colors)]
            ser.marker.format.line.color.rgb = WHITE
            if label_series is not None and i == label_series:
                dl = ser.data_labels
                dl.show_value = True
                dl.number_format = '#,##0'
                dl.number_format_is_linked = False
                dl.font.size = Pt(7.5)
                dl.font.color.rgb = INK
                dl.position = XL_LABEL_POSITION.ABOVE
        return gfx

    # ----------------------------------------------------------- closing ----
    def verify(self):
        """Raise if the footer denominator does not match the real count.

        The June cycle shipped wrong footers twice. This is why n_slides is
        declared up front and checked before save, not inferred at the end.
        """
        actual = len(self.prs.slides)
        if actual != self.n_slides:
            raise AssertionError(
                f'footer says "/ {self.n_slides}" but the deck has {actual} '
                f'slides. Fix n_slides or the slide list, then rebuild.')
        return actual

    def save(self, path):
        self.verify()
        os.makedirs(os.path.dirname(os.path.abspath(path)), exist_ok=True)
        self.prs.save(path)
        return path


def _smoke():
    """Build a deck exercising every primitive. The runnable check."""
    d = Deck(entity='SMOKE TEST', month='Month Year', n_slides=3)
    d.cover('QUARTERLY BUSINESS REVIEW', 'Entity - Paid Media',
            'Month Year Performance', 'Channel . Channel . Channel',
            'Reporting period: 1-31 | Currency: USD | Prepared by Berelvant')

    s = d.content('EXECUTIVE SUMMARY', 'One headline, navy Georgia, one line.')
    d.table(s, 0.5, 1.9, 7.4, 2.0, [
        ['Channel', 'Spend', 'Impr', 'Clicks', 'Conv', 'CPA'],
        ['Search', '$10,624', '90,394', '4,090', '20', '$531'],
        ['Social', '$25,924', '19.78M', '407,136', 'Held*', '—'],
        ['Total', '$36,548', '19.87M', '411,226', '—', '—'],
    ], col_widths=[2.33, 1.06, 1.06, 1.06, 0.85, 1.06], total_last=True)
    d.note(s, 0.5, 4.02, 7.4, 0.45,
           'Conversions come from different systems and funnel stages, so no '
           'blended total is shown.')
    d.card(s, 8.1, 1.9, 4.7, 2.0, 'What the data shows',
           ['Green edge carries the data semantics.'], GREEN)
    d.card(s, 8.1, 4.1, 4.7, 2.0, 'What we recommend',
           ['Gold edge carries the recommendation semantics.'], GOLD)
    # tiles stay inside the 7.4in left column so they never overlap the
    # cards at x=8.1 (the first render of this smoke test did exactly that,
    # and the tile's white fill clipped a card's accent edge and first glyph)
    for i, (lbl, val) in enumerate(
            [('Spend', '$36,548'), ('Clicks', '411,226'),
             ('Submitted apps', '20'), ('Cost / app', '$531')]):
        d.tile(s, 0.5 + i * 1.87, 4.6, lbl, val, sub='vs prior month', w=1.75)

    s = d.content('RECOMMENDATIONS AND BLOCKERS', 'Four cards, then the strip.')
    for i, (lbl, col) in enumerate([('Critical', CORAL), ('High', GOLD),
                                    ('Opportunity', GREEN),
                                    ('Measurement to close', DEEP)]):
        d.blocker(s, 0.5 + i * 3.15, 1.9, 2.95, 2.4, lbl,
                  ['First item.', 'Second item.'], col)
    d.line_chart(s, 0.5, 4.5, 6.0, 2.4, ['Jan', 'Feb', 'Mar', 'Apr'],
                 [('Sessions', [9380, 22229, 12614, 7577])], label_series=0)
    d.strip(s, 6.8, 4.5, 6.03, 0.8, 'One white Georgia takeaway line.')

    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       '..', '_scratch', 'housestyle-smoke.pptx')
    d.save(out)

    # the check: every primitive ran, and the footer denominator is honest
    assert d.verify() == 3
    bad = Deck('X', 'Y', n_slides=9)
    bad.content('A', 'B')
    try:
        bad.verify()
        raise SystemExit('FAIL: verify() did not catch a wrong slide count')
    except AssertionError:
        pass
    print(f'OK  3 slides, every primitive rendered, verify() catches drift\n'
          f'    {os.path.abspath(out)}')


if __name__ == '__main__':
    _smoke()
