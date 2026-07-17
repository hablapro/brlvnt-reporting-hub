#!/usr/bin/env python3
"""Restyle June 2026 FOREX decks (GGMI + GCG) to the forex-reporting-deck house style.

Content, geometry, and slide order are preserved. Only the visual system changes:
palette, fonts, band/rule/logo treatment, card borders, footers, table skins.
"""
import copy, sys
from pptx import Presentation
from pptx.util import Emu, Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

A = '{http://schemas.openxmlformats.org/drawingml/2006/main}'
R = '{http://schemas.openxmlformats.org/officeDocument/2006/relationships}'

# ---- house tokens -----------------------------------------------------------
DEEP    = RGBColor(0x0F, 0x15, 0x35)
NAVY    = RGBColor(0x1E, 0x27, 0x61)
CORAL   = RGBColor(0xF9, 0x61, 0x67)
GOLD    = RGBColor(0xE9, 0xB4, 0x4C)
GREEN   = RGBColor(0x2E, 0x7D, 0x57)
ICE     = RGBColor(0xCA, 0xDC, 0xFC)
MUTED   = RGBColor(0x5A, 0x60, 0x72)
BORDER  = RGBColor(0xD0, 0xD4, 0xDC)
LIGHT   = RGBColor(0xF9, 0xFA, 0xFC)
INK     = RGBColor(0x2B, 0x31, 0x47)
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)

# ---- old palette ------------------------------------------------------------
O_NAVY  = RGBColor(0x1A, 0x2A, 0x3A)
O_BLUE  = RGBColor(0x3B, 0x59, 0x98)
O_GRAY  = RGBColor(0xF3, 0xF3, 0xF3)
O_GREEN = RGBColor(0x27, 0xAE, 0x60)
O_RED   = RGBColor(0xE7, 0x4C, 0x3C)
O_BODY  = RGBColor(0x32, 0x37, 0x3E)
O_MUT   = RGBColor(0x78, 0x82, 0x8C)
O_SUB   = RGBColor(0x5A, 0x5A, 0x5A)

TEXTMAP = {O_BLUE: NAVY, O_NAVY: INK, O_BODY: INK, O_GREEN: GREEN,
           O_RED: CORAL, O_MUT: MUTED, O_SUB: MUTED}

BERELVANT_LOGO = "forex-skill/forex-reporting-deck/assets/berelvant_logo.png"

GOLD_KEYS  = ("RECOMMEND", "SCALING", "OPPORTUNITY", "NEXT STEP", "WHAT WE DO")
GREEN_KEYS = ("WHAT IT MEANS", "HIGHLIGHTS", "WHAT THE DATA", "DATA SHOWS")
CORAL_KEYS = ("HONEST READ", "PRIORITY", "CRITICAL", "READ", "RISK")
DEEP_KEYS  = ("MEASUREMENT",)


def inch(v):
    return Emu(int(v)).inches


def solid_rgb(sh):
    try:
        if sh.fill.type == 1:
            return sh.fill.fore_color.rgb
    except Exception:
        pass
    return None


def set_border(sh, color, w_pt):
    sh.line.color.rgb = color
    sh.line.width = Pt(w_pt)


def iter_runs(tf):
    for p in tf.paragraphs:
        for r in p.runs:
            yield p, r


def style_run(r, name=None, color=None, bold=None, size=None):
    if name is not None:
        r.font.name = name
    if color is not None:
        r.font.color.rgb = color
    if bold is not None:
        r.font.bold = bold
    if size is not None:
        r.font.size = Pt(size)


def first_text(sh):
    if not sh.has_text_frame:
        return ""
    return sh.text_frame.text.strip()


def label_color(text):
    t = text.upper()
    for k in DEEP_KEYS:
        if k in t:
            return DEEP
    for k in GREEN_KEYS:
        if k in t:
            return GREEN
    for k in GOLD_KEYS:
        if k in t:
            return GOLD
    for k in CORAL_KEYS:
        if k in t:
            return CORAL
    return None


def remap_text(sh, in_band=False, on_dark=False):
    """Font + color remap for one text shape."""
    if not sh.has_text_frame:
        return
    for p, r in iter_runs(sh.text_frame):
        # font family
        sz = r.font.size.pt if r.font.size else None
        cur = None
        try:
            if r.font.color and r.font.color.type is not None:
                cur = r.font.color.rgb
        except Exception:
            cur = None
        header_like = in_band or (sz is not None and sz >= 15)
        style_run(r, name=('Georgia' if header_like else 'Calibri'))
        if on_dark or in_band:
            continue  # keep white / light text on dark shapes
        if cur in TEXTMAP:
            new = TEXTMAP[cur]
            # big headlines go navy, not ink
            if header_like and new == INK:
                new = NAVY
            style_run(r, color=new)


def swap_picture(pic, img_path):
    """Point an existing picture shape at a new image file, keeping geometry."""
    part = pic.part
    image_part, rId = part.get_or_add_image_part(img_path)
    blip = pic._element.blipFill.find(A + 'blip')
    blip.set(R + 'embed', rId)


def add_rect(slide, x, y, w, h, fill, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Emu(int(x)), Emu(int(y)), Emu(int(w)), Emu(int(h)))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def table_borders(tbl):
    for row in tbl.rows:
        for cell in row.cells:
            tcPr = cell._tc.get_or_add_tcPr()
            for tag in ('lnL', 'lnR', 'lnT', 'lnB'):
                for old in tcPr.findall(A + tag):
                    tcPr.remove(old)
            for tag in ('lnL', 'lnR', 'lnT', 'lnB'):
                ln = etree.SubElement(tcPr, A + tag)
                ln.set('w', '6350')  # 0.5pt
                ln.set('cap', 'flat')
                fill = etree.SubElement(ln, A + 'solidFill')
                clr = etree.SubElement(fill, A + 'srgbClr')
                clr.set('val', 'D0D4DC')
                # order: solidFill first is fine
            # move ln* elements before other children order requirements are lax for renderers
    return


def restyle_table(gfx):
    tbl = gfx.table
    n = len(tbl.rows)
    for ri, row in enumerate(tbl.rows):
        for ci, cell in enumerate(row.cells):
            rgb = None
            try:
                if cell.fill.type == 1:
                    rgb = cell.fill.fore_color.rgb
            except Exception:
                rgb = None
            first_cell_text = row.cells[0].text_frame.text.strip().lower()
            is_total = first_cell_text.startswith(('total', 'q3 total', 'combined'))
            if ri == 0 or rgb == O_BLUE:
                cell.fill.solid(); cell.fill.fore_color.rgb = NAVY
            elif is_total:
                cell.fill.solid(); cell.fill.fore_color.rgb = ICE
            elif rgb == O_GRAY:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    style_run(r, name='Calibri')
                    if ri == 0 or rgb == O_BLUE:
                        style_run(r, color=WHITE)
                    elif is_total:
                        style_run(r, color=NAVY, bold=True)
                    else:
                        cur = None
                        try:
                            if r.font.color and r.font.color.type is not None:
                                cur = r.font.color.rgb
                        except Exception:
                            pass
                        if cur in TEXTMAP:
                            style_run(r, color=TEXTMAP[cur])
    table_borders(tbl)


def send_to_back(slide, sh):
    spTree = slide.shapes._spTree
    spTree.remove(sh._element)
    spTree.insert(2, sh._element)


def restyle_cover(slide, W, H, entity_label):
    bg = add_rect(slide, 0, 0, W, H, DEEP)
    send_to_back(slide, bg)
    for sh in list(slide.shapes):
        rgb = solid_rgb(sh)
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and rgb is not None:
            x, y, w, h = sh.left, sh.top, sh.width, sh.height
            if inch(w) < 0.15 and inch(h) > 0.5:
                sh.fill.fore_color.rgb = CORAL          # vertical accent bar
            elif inch(w) > 12 and inch(h) > 5:
                sh.fill.fore_color.rgb = DEEP           # full-bleed background rect
            elif rgb in (O_NAVY, DEEP):
                sh.fill.fore_color.rgb = DEEP           # bottom strip blends in
            continue
        if sh.has_text_frame:
            txt = first_text(sh).upper()
            for p, r in iter_runs(sh.text_frame):
                sz = r.font.size.pt if r.font.size else 12
                if 'QUARTERLY BUSINESS REVIEW' in txt:
                    style_run(r, name='Georgia', color=CORAL, bold=True)
                elif sz >= 28:
                    style_run(r, name='Georgia', color=WHITE, bold=True)
                elif sz >= 16:
                    style_run(r, name='Georgia', color=ICE)
                else:
                    style_run(r, name='Calibri', color=ICE)
    # Berelvant logo bottom-right
    slide.shapes.add_picture(BERELVANT_LOGO, W - Inches(2.5), H - Inches(0.95), Inches(2.0), Inches(0.57))


def restyle_content(slide, W, H, entity_label, month_label, idx, total):
    band = None
    for sh in slide.shapes:
        if sh.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE and solid_rgb(sh) == O_NAVY \
           and inch(sh.top) < 0.3 and inch(sh.width) > 10:
            band = sh
            break

    thin_h = []   # thin horizontal bars (card top-bars / strip accents)
    bodies = []   # gray card bodies
    strips = []   # closing takeaway strips (deep navy)
    for sh in list(slide.shapes):
        st = sh.shape_type
        rgb = solid_rgb(sh)
        if getattr(sh, 'has_table', False) and sh.has_table:
            restyle_table(sh)
            continue
        if st == MSO_SHAPE_TYPE.PICTURE:
            if inch(sh.top) < 0.45 and inch(sh.left) > 11:
                swap_picture(sh, BERELVANT_LOGO)        # band logo -> Berelvant
            continue
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE and rgb is not None:
            w, h = inch(sh.width), inch(sh.height)
            if sh is band:
                sh.fill.fore_color.rgb = NAVY
                remap_text(sh, in_band=True)
                continue
            if rgb in (O_NAVY,) and inch(sh.top) > H.inches * 0.55 and w > 9:
                # closing takeaway strip body
                sh.fill.fore_color.rgb = DEEP
                for p, r in iter_runs(sh.text_frame):
                    style_run(r, name='Georgia', color=WHITE, bold=True)
                continue
            if rgb == O_GRAY:
                if inch(sh.top) >= H.inches - 0.45:
                    sh.fill.fore_color.rgb = LIGHT      # footer band stays light
                    continue
                if w > 9 and inch(sh.top) > H.inches * 0.6:
                    # closing takeaway strip: deep navy, white Georgia
                    sh.fill.fore_color.rgb = DEEP
                    for p, r in iter_runs(sh.text_frame):
                        style_run(r, name='Georgia', color=WHITE, bold=True)
                    strips.append(sh)
                    continue
                bodies.append(sh)
                continue
            if h <= 0.14 and w >= 0.5:
                thin_h.append(sh)
                continue
            if w <= 0.14 and h >= 0.3:
                sh.fill.fore_color.rgb = CORAL          # stat-card left edge
                continue
            if rgb == O_GREEN:
                sh.fill.fore_color.rgb = GREEN
            elif rgb == O_RED:
                sh.fill.fore_color.rgb = CORAL
            elif rgb == O_BLUE:
                sh.fill.fore_color.rgb = NAVY
            elif rgb == O_NAVY:
                sh.fill.fore_color.rgb = NAVY
                remap_text(sh, on_dark=True)
                continue
        remap_text(sh)

    # card bodies: white + border; label decides the card's accent color
    card_info = []  # (body, label_shape, accent)
    for body in bodies:
        label, label_sh = "", None
        for sh in slide.shapes:
            if sh.has_text_frame and sh.top >= body.top - Emu(50000) \
               and sh.top <= body.top + body.height and abs(sh.left - body.left) < Inches(0.6):
                t = first_text(sh)
                if t and t[:40].upper() == t[:40] and any(c.isalpha() for c in t):
                    label, label_sh = t, sh
                    break
        key = (label or first_text(body)).upper()
        accent = label_color(key) or NAVY
        body.fill.solid()
        body.fill.fore_color.rgb = WHITE
        if accent is CORAL and 'READ' in key:
            set_border(body, CORAL, 1.25)               # house read box
        else:
            set_border(body, BORDER, 0.75)
        remap_text(body)
        card_info.append((body, label_sh, accent))
        # label text: read boxes take a navy Georgia title (house spec);
        # other cards take their accent color (navy when accent is deep navy)
        if label_sh is not None:
            if accent is CORAL and 'READ' in key:
                lab_rgb = NAVY
            else:
                lab_rgb = NAVY if accent in (DEEP, NAVY) else accent
            for r in label_sh.text_frame.paragraphs[0].runs:
                style_run(r, name='Georgia', color=lab_rgb, bold=True)

    for bar in thin_h:
        rgb = solid_rgb(bar)
        placed = False
        for body, label_sh, accent in card_info:
            near_left = abs(inch(body.left) - inch(bar.left)) < 0.3
            at_top = 0 <= inch(body.top) - inch(bar.top) <= 0.3
            at_bottom = abs(inch(bar.top) - (inch(body.top) + inch(body.height))) < 0.1
            inside_bottom = abs((inch(bar.top) + inch(bar.height)) - (inch(body.top) + inch(body.height))) < 0.1
            if near_left and (at_top or at_bottom or inside_bottom):
                bar.fill.fore_color.rgb = CORAL if (accent is CORAL) else accent
                placed = True
                break
        if placed:
            continue
        near_strip = any(abs(inch(bar.top) + inch(bar.height) - inch(s.top)) < 0.12
                         and abs(inch(bar.left) - inch(s.left)) < 0.3 for s in strips)
        if near_strip:
            bar.fill.fore_color.rgb = CORAL
        elif rgb == O_GREEN:
            bar.fill.fore_color.rgb = GREEN
        elif rgb == O_RED:
            bar.fill.fore_color.rgb = CORAL
        elif rgb in (O_BLUE, O_NAVY):
            bar.fill.fore_color.rgb = NAVY

    # text overlapping a deep-navy takeaway strip goes white Georgia bold
    for strip in strips:
        for sh in slide.shapes:
            if sh is strip or not sh.has_text_frame:
                continue
            if sh.top >= strip.top - Emu(30000) and sh.top < strip.top + strip.height \
               and sh.left >= strip.left - Emu(30000):
                for p, r in iter_runs(sh.text_frame):
                    style_run(r, name='Georgia', color=WHITE, bold=True)

    # band-less slides (user-added QBR slides): headline goes house navy
    if band is None:
        for sh in slide.shapes:
            if sh.has_text_frame and inch(sh.top) < 1.0:
                for p, r in iter_runs(sh.text_frame):
                    sz = r.font.size.pt if r.font.size else None
                    if sz is not None and sz >= 15:
                        cur = None
                        try:
                            if r.font.color and r.font.color.type is not None:
                                cur = r.font.color.rgb
                        except Exception:
                            pass
                        if cur != WHITE:
                            style_run(r, name='Georgia', color=NAVY, bold=True)

    # coral rule beneath the band
    if band is not None:
        add_rect(slide, 0, band.top + band.height, W, Inches(0.05), CORAL)

    # footer label
    for sh in slide.shapes:
        if sh.has_text_frame and inch(sh.top) > H.inches - 0.6:
            t = first_text(sh)
            if t.lower().startswith('performance review') or t == 'Performance Review':
                sh.text_frame.paragraphs[0].runs[0].text = \
                    f"{entity_label}   |   {month_label} Performance Review   |   Berelvant"
            for p, r in iter_runs(sh.text_frame):
                style_run(r, name='Calibri', color=MUTED)


def restyle(path_in, path_out, entity_label, month_label):
    prs = Presentation(path_in)
    W, H = Emu(prs.slide_width), Emu(prs.slide_height)
    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, 1):
        if i == 1:
            restyle_cover(slide, W, H, entity_label)
        else:
            restyle_content(slide, W, H, entity_label, month_label, i, total)
    prs.save(path_out)
    print(f"saved {path_out}")


if __name__ == '__main__':
    base = "/Users/rpro/AI-BRLVNT/Brlvnt-Reporting-Analytics-2026/report-client-decks"
    restyle(f"{base}/06. GGMI_LATAM_June_2026_Performance_Review.pptx",
            "render-after/ggmi-house.pptx", "GGMI (LATAM)", "June 2026")
    restyle(f"{base}/06. GCG_US_June_2026_Performance_Review.pptx",
            "render-after/gcg-house.pptx", "GCG (US Hispanic)", "June 2026")
