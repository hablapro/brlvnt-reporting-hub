#!/usr/bin/env python3
"""Build the GGMI June 2026 Performance Review PPTX working draft.

Mirrors the structure of the FINAL May client deck (Google Slides
1npxoxCC..., 11 slides: per-channel Performance + Read/Recommendations
pairs) using the visual system of the May PPTX draft
(report-client-decks/05. GGMI_LATAM_May_2026_Performance_Review.pptx).
June adds one Site Traffic (GA4) slide -> 12 slides.

Numbers source: reports/forex/ggmi/2026-06/model/ + data/ workbooks.
Narrative source: output/GGMI-Jun-2026-narrative-draft.md (v2, approved
by Renzo 2026-07-16). Client-facing rules: Meta conversions held, raw
Azerion spend, no vendor/competitor names in commentary.
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION,
                             XL_MARKER_STYLE, XL_LABEL_POSITION)

NAVY = RGBColor(0x1A, 0x2A, 0x3A)
BLUE = RGBColor(0x3B, 0x59, 0x98)
GRAY_BG = RGBColor(0xF3, 0xF3, 0xF3)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
BODY = RGBColor(0x32, 0x37, 0x3E)
MUTED = RGBColor(0x78, 0x82, 0x8C)
SUB = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, 'assets', 'logo_1_fd6b9b.png')
OUT = os.path.join(HERE, '..', '..', 'report-client-decks',
                   '06. GGMI_LATAM_June_2026_Performance_Review.pptx')
BREADCRUMB = 'FOREX.com  |  GGMI (LATAM)  ·  June 2026'
FOOTER_LEFT = 'Performance Review'
N_SLIDES = 17
CHART_BLUE = RGBColor(0x3B, 0x59, 0x98)
CHART_GREEN = RGBColor(0x27, 0xAE, 0x60)
GRID = RGBColor(0xD9, 0xD9, 0xD9)
MONTHS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill, line=False):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if not line:
        sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=9.5, bold=False, color=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    """runs: str, or list of paragraph strings."""
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = [runs] if isinstance(runs, str) else runs
    for i, ptext in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0:
            p.space_before = Pt(6)
        r = p.add_run()
        r.text = ptext
        r.font.name = 'Arial'
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(slide, section, headline, pageno):
    rect(slide, 0, 0, 13.333, 0.5, NAVY)
    text(slide, 0.5, 0.1, 9.0, 0.3, BREADCRUMB, 9, True, WHITE)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.2), Inches(0.11),
                                 Inches(0.98), Inches(0.28))
    text(slide, 0.5, 0.62, 12.0, 0.3, section, 11, True, BLUE)
    text(slide, 0.5, 0.95, 12.3, 0.75, headline, 20, True, NAVY)
    rect(slide, 0, 7.2, 13.333, 0.3, GRAY_BG)
    text(slide, 0.5, 7.23, 6.0, 0.22, FOOTER_LEFT, 8, False, MUTED)
    text(slide, 12.1, 7.23, 0.9, 0.22, f'{pageno} / {N_SLIDES}', 8, False,
         MUTED, align=PP_ALIGN.RIGHT)


def tile(slide, x, y, label, value, sub=None, w=1.9):
    h = 1.1 if sub else 1.0
    rect(slide, x, y, w, h, GRAY_BG)
    rect(slide, x, y, 0.1, h, BLUE)
    text(slide, x + 0.2, y + 0.05, w - 0.3, 0.3, label, 8, True, MUTED)
    text(slide, x + 0.2, y + 0.32, w - 0.3, 0.5, value, 15, True, NAVY)
    if sub:
        text(slide, x + 0.2, y + 0.82, w - 0.3, 0.24, sub, 8, False, SUB)


def card(slide, x, y, w, h, title, body_paras, accent=BLUE):
    rect(slide, x, y, w, h, GRAY_BG)
    rect(slide, x, y, w, 0.1, accent)
    text(slide, x + 0.2, y + 0.2, w - 0.4, 0.3, title, 10, True, accent)
    text(slide, x + 0.2, y + 0.55, w - 0.4, h - 0.75, body_paras, 9.5,
         False, BODY)


def line_chart(slide, x, y, w, h, series, colors, label_series=None):
    """series: list of (name, values); label_series: index to direct-label."""
    cd = CategoryChartData()
    cd.categories = MONTHS
    for name, vals in series:
        cd.add_series(name, vals)
    gfx = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x),
                                 Inches(y), Inches(w), Inches(h), cd)
    ch = gfx.chart
    ch.has_title = False
    ch.font.size = Pt(8)
    ch.font.name = 'Arial'
    ch.font.color.rgb = BODY
    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(8)
    val_ax = ch.value_axis
    val_ax.has_major_gridlines = True
    val_ax.major_gridlines.format.line.color.rgb = GRID
    val_ax.major_gridlines.format.line.width = Pt(0.5)
    val_ax.tick_labels.number_format = '#,##0'
    val_ax.tick_labels.number_format_is_linked = False
    val_ax.format.line.fill.background()
    cat_ax = ch.category_axis
    cat_ax.has_major_gridlines = False
    cat_ax.format.line.color.rgb = GRID
    for i, s in enumerate(ch.series):
        s.smooth = False
        s.format.line.color.rgb = colors[i]
        s.format.line.width = Pt(2.25)
        s.marker.style = XL_MARKER_STYLE.CIRCLE
        s.marker.size = 6
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = colors[i]
        s.marker.format.line.color.rgb = WHITE
        if label_series is not None and i == label_series:
            s.data_labels.show_value = True
            s.data_labels.number_format = '#,##0'
            s.data_labels.number_format_is_linked = False
            s.data_labels.font.size = Pt(7.5)
            s.data_labels.font.color.rgb = BODY
            s.data_labels.position = XL_LABEL_POSITION.ABOVE
    return gfx


def table(slide, x, y, w, h, rows, col_widths=None, bold_last=False):
    n_r, n_c = len(rows), len(rows[0])
    gfx = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = BLUE
            elif bold_last and ri == n_r - 1:
                cell.fill.fore_color.rgb = GRAY_BG
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else GRAY_BG
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.name = 'Arial'
            r.font.size = Pt(8.5)
            r.font.bold = (ri == 0) or (ci == 0) or (bold_last and ri == n_r - 1)
            r.font.color.rgb = WHITE if ri == 0 else BODY
            if ci > 0:
                p.alignment = PP_ALIGN.RIGHT
    return gfx


# ---- Slide 1: Title -------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, WHITE)
rect(s, 0, 6.4, 13.333, 1.1, NAVY)
rect(s, 0.5, 2.2, 0.1, 2.5, BLUE)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.5), Inches(0.5), Inches(2.5), Inches(0.71))
text(s, 0.9, 2.2, 11.0, 0.4, 'QUARTERLY BUSINESS REVIEW', 14, True, BLUE)
text(s, 0.9, 2.7, 12.0, 1.0, 'GGMI (LATAM) — Paid Media', 36, True, NAVY)
text(s, 0.9, 3.8, 12.0, 0.6,
     'Q3 FY2026 Business Review  ·  June 2026 Performance', 22, False, NAVY)
text(s, 0.9, 4.5, 12.0, 0.4, 'Bing Ads · Meta · Azerion · Quantcast',
     12, False, SUB)
text(s, 0.9, 6.65, 12.0, 0.4,
     'Reporting period: June 1–30, 2026   |   Currency: USD   |   '
     'Prepared by Berelvant · 07.16.2026', 10, False, WHITE)

# ---- Slide 2: Executive Summary -------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'EXECUTIVE SUMMARY',
       'The biggest media month of 2026. Efficiency held where we can '
       'measure it.', 8)
table(s, 0.5, 1.9, 7.4, 2.0, [
    ['Channel', 'Spend', 'Impr', 'Clicks', 'Conv', 'CPA'],
    ['Bing Ads', '$25,659', '466,582', '21,480', '50', '$513'],
    ['Meta', '$25,924', '19.78M', '407,136', 'Held*', '—'],
    ['Azerion', '$35,026', '7.68M', '9,910', '42', '$834'],
    ['Quantcast', '$33,784', '41.96M', '11,284', '11 VT*', '—'],
    ['Total', '$120,393', '69.9M', '449,810', '—', '—'],
], col_widths=[2.33, 1.06, 1.06, 1.06, 0.85, 1.06], bold_last=True)
text(s, 0.5, 4.02, 7.4, 0.45,
     '* Paid-social conversions are held pending a placement audit. '
     'Programmatic view-through is directional only. Conversions use '
     'different events per channel and are never summed.', 7.5, False, SUB)
card(s, 8.1, 1.9, 4.7, 2.4, 'HEADLINES', [
    'June was the biggest GGMI media month of 2026: $120,393 across four '
    'channels, up 53% from May. Search produced 50 submitted applications '
    'at a $513 CPA, up from 33 in May, and Azerion held efficiency with '
    '42 applications through a 20% budget increase.',
], GREEN)
card(s, 0.5, 4.6, 7.4, 2.4, 'WHAT THE DATA SAYS', [
    'Site traffic recovered with the spend: 9,236 Mexico sessions, up 61% '
    'from May, and 5,838 unique visitors, up 125%. The recovery is real '
    'but bought. Media drove essentially all of the June gain while '
    'organic search sat at its low.',
    'We are approximately 31% into the annual GGMI budget, so H2 has room '
    'to scale once the quality fixes land.',
], BLUE)
card(s, 8.1, 4.6, 4.7, 2.4, 'JUNE PRIORITY', [
    'Concentrate search spend in Mexico (49% served out-of-market in '
    'June), audit paid-social placements before reinstating its '
    'conversion reporting, and open the SEO workstream for the organic '
    'decline.',
], RED)

# ---- Slide 3: Cont' Executive Summary --------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "CONT' EXECUTIVE SUMMARY",
       'Three findings shape the July plan.', 9)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT IS STILL FORMING', [
    'The June traffic recovery was bought by media. Organic search has '
    'fallen 75% since January, confirmed independently in Search Console, '
    'for identifiable SEO reasons. No media budget fixes this, and it is '
    'why total site traffic feels flat while investment grows.',
    'Meta scaled spend 291% and delivery stayed cheap, but click quality '
    'dropped sharply as it scaled. Its conversion reporting stays on hold '
    'until the placement audit completes.',
    '49% of June search spend served users outside Mexico. The fix is '
    'defined and partially applied; completing it is the top efficiency '
    'lever in the account.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Complete the Mexico-only restriction on search. Roughly $12.6K per '
    'month is at stake and the June 3 exclusion is already working.',
    'Run the placement-level audit and exclusions on paid social before '
    'scaling it further, and judge the channel on measured site sessions '
    'until validation completes.',
    'Open a dedicated SEO workstream for the Spanish site. It is the '
    'structural fix for traffic that tracks media spend instead of '
    'compounding on top of it.',
    'Apply the 49-site programmatic blocklist and set a campaign-level '
    'viewability floor.',
], GREEN)

# ---- Slide 4: Bing Performance ---------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'BING ADS  ·  PERFORMANCE',
       '50 submitted applications at $513. Efficiency held through a '
       '61% scale-up.', 10)
tile(s, 0.5, 1.9, 'SPEND', '$25,659', '+61% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '466,582', '3 campaigns')
tile(s, 4.5, 1.9, 'CLICKS', '21,480', 'CTR 4.60%')
tile(s, 6.5, 1.9, 'AVG CPC', '$1.19', 'May: $0.53')
tile(s, 8.5, 1.9, 'SUBMITTED APPS', '50', 'May: 33')
tile(s, 10.5, 1.9, 'CPA', '$513', 'May: $484')
text(s, 0.5, 3.35, 7.0, 0.3, 'Campaigns', 11, True, NAVY)
table(s, 0.5, 3.65, 7.6, 1.4, [
    ['Campaign', 'Spend', 'Clicks', 'Apps', 'CPA'],
    ['AO Generic (policy test)', '$13,508', '11,020', '27', '$500'],
    ['Brand + Generic', '$8,154', '7,559', '9', '$906'],
    ['Platform Intercept', '$3,996', '2,901', '14', '$285'],
], col_widths=[3.0, 1.2, 1.2, 1.0, 1.2])
card(s, 8.4, 3.35, 4.4, 3.55, 'HIGHLIGHTS', [
    'Conversions rose 52% (33 to 50) while spend rose 61%, so the CPA '
    'moved only $484 to $513 through the scale-up.',
    'The account expanded from 1 campaign to 3. Platform Intercept is the '
    'efficiency leader at a $285 CPA.',
    'CPC rose $0.53 to $1.19 with the shift into the new campaign themes.',
], GREEN)
text(s, 0.5, 5.25, 7.6, 1.6, [
    '51% of June spend served Mexico. The remaining 49% ($12,637) served '
    'out-of-market users. The June 3 exclusion is already working, and '
    'completing the Mexico-only restriction is optimization #1 for July.',
], 9.5, False, BODY)

# ---- Slide 5: Bing Read & Recommendations ----------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'BING ADS  ·  READ & RECOMMENDATIONS',
       'Concentrate every search dollar in Mexico.', 11)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT THE DATA SHOWS', [
    'Search remains the most efficient measured converter in the account, '
    'and it held that efficiency while scaling 61%.',
    'Only 51% of June spend served Mexico. Out-of-market paid visitors '
    'outnumbered Mexican paid visitors on the site for most of H1, so '
    'concentrating that spend in-market should improve effective Mexico '
    'CPA immediately.',
    'The new campaign themes raised CPC ($0.53 to $1.19). Platform '
    'Intercept converts at $285, so the mix has room to rebalance toward '
    'what works.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Complete the Mexico presence-only restriction. It is the top '
    'efficiency lever in the account at roughly $12.6K per month.',
    'Rebalance budget toward the proven converters, led by Platform '
    'Intercept at $285, while the newer themes prove out.',
    'Link the search account in GA4 so July reporting attributes paid '
    'search fully (detail on the measurement slide).',
], GREEN)

# ---- Slide 6: Meta Performance ---------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'META  ·  PERFORMANCE',
       'Reach scaled 4x at low cost. Click quality is the June question.', 12)
tile(s, 0.5, 1.9, 'SPEND', '$25,924', '+291% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '19.78M', '100% Mexico')
tile(s, 4.5, 1.9, 'LINK CLICKS', '407,136', 'CTR 2.06%')
tile(s, 6.5, 1.9, 'CPM', '$1.31', 'May: $0.97')
tile(s, 8.5, 1.9, 'LP VIEWS', '249,972', 'platform-reported')
tile(s, 10.5, 1.9, 'CONVERSIONS', 'Held', 'validation open')
card(s, 0.5, 3.35, 6.1, 3.55, 'HIGHLIGHTS', [
    'Delivery scaled hard and stayed cheap: 19.8M impressions at a $1.31 '
    'CPM, all in Mexico, with the first month off a pure traffic '
    'objective.',
    'This account converted platform clicks into measured site visits at '
    '15.7% in February. June ran at 0.19%. Same pixel, same site, same '
    'market. The difference is the inventory the June campaigns bought.',
], GREEN)
card(s, 6.9, 3.35, 5.9, 3.55, 'WHY CONVERSIONS ARE HELD', [
    'The platform reports 86 conversions for June. 67 of them come from '
    'one retargeting campaign with only 126 landing-page views, which is '
    'not a credible application rate.',
    'We are holding conversion reporting until the placement audit and '
    'pixel-event review complete. Delivery metrics (reach, clicks, CPM) '
    'are unaffected.',
], RED)

# ---- Slide 7: Meta Read & Recommendations ----------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'META  ·  READ & RECOMMENDATIONS',
       'Audit the placements before judging the channel.', 13)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT THE DATA SHOWS', [
    'June buying skewed to low-cost in-app inventory, and 63% of spend '
    'served users 55 and older in a market where the target skews '
    'younger.',
    'Platform-reported landing-page views exceed what any measurement '
    'system records on the site by an order of magnitude. The landing '
    'pages are confirmed tagged and tracking, so the February-vs-June '
    'capture comparison isolates inventory quality, not tracking, as the '
    'cause.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Run the placement-level audit and apply exclusions. This is the '
    'highest-value action on this channel for July.',
    'Rebalance the 55+ age skew and complete the shift to a conversion '
    'objective.',
    'Judge the channel on measured site sessions until conversion '
    'validation completes, then reinstate conversion reporting.',
], GREEN)

# ---- Slide 8: Azerion Performance ------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'AZERION  ·  PERFORMANCE',
       '42 submitted applications. Efficiency roughly held at higher '
       'volume.', 14)
tile(s, 0.5, 1.9, 'SPEND', '$35,026', '+20% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '7.68M', '2.18M devices')
tile(s, 4.5, 1.9, 'CLICKS', '9,910', '+57% MoM')
tile(s, 6.5, 1.9, 'STEP 1', '440', 'May: 225')
tile(s, 8.5, 1.9, 'APPLICATIONS', '42', '$834 CPA')
tile(s, 10.5, 1.9, 'VIEWABILITY', '68.5%', '70% standard')
text(s, 0.5, 3.35, 6.0, 0.3, 'Ad Sets (Applications)', 11, True, NAVY)
table(s, 0.5, 3.65, 6.0, 2.6, [
    ['Ad Set', 'Spend', 'Apps', 'CPA'],
    ['Experience', '$9,118', '11', '$829'],
    ['Global Market', '$4,124', '10', '$412'],
    ['TradeForex', '$4,221', '7', '$603'],
    ['Language Broker', '$4,994', '7', '$713'],
    ['Conversion_Instruments', '$4,211', '4', '$1,053'],
    ['Spanish Platform', '$4,229', '3', '$1,410'],
    ['Conversion_Commodities', '$4,128', '0', '—'],
], col_widths=[2.8, 1.1, 0.9, 1.2])
card(s, 6.8, 3.35, 6.0, 3.55, 'HIGHLIGHTS', [
    'Applications rose 14% (37 to 42) on a 20% budget increase. The CPA '
    'moved $792 to $834, so efficiency roughly held at higher volume.',
    'Application starts nearly doubled (225 to 440), keeping the upper '
    'funnel active.',
    'Global Market converted at $412 and Experience carried the most '
    'volume. Conversion_Commodities spent $4,128 with no applications '
    'and is the first reallocation candidate.',
], GREEN)

# ---- Slide 9: Azerion Read & Recommendations -------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'AZERION  ·  READ & RECOMMENDATIONS',
       'Hold the pace. Confirm the delivery detail.', 15)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT THE DATA SHOWS', [
    'Second consecutive month of scaled application volume, with '
    'efficiency holding within 6% while budget grew 20%.',
    'Viewability is 68.5%, just under the 70% standard.',
    'We have asked the vendor for country-level delivery confirmation '
    'plus funnel, site, format, and creative breakdowns (requested early '
    'July; awaiting reply). Mexico-only delivery certification is pending '
    'that response.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Hold the budget pace until the delivery breakdowns arrive.',
    'Concentrate budget on the converting ad sets, led by Global Market '
    '($412) and Experience, and reallocate away from '
    'Conversion_Commodities ($4,128, zero applications).',
    'Push viewability over the 70% standard.',
], GREEN)

# ---- Slide 10: Quantcast ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'QUANTCAST  ·  PERFORMANCE & READ',
       'A deliberate reach month. Cheap delivery, with a quality bill to '
       'manage.', 16)
tile(s, 0.5, 1.9, 'SPEND', '$33,784', '+26% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '41.96M', '+155% MoM')
tile(s, 4.5, 1.9, 'CPM', '$0.81', 'May: $1.63')
tile(s, 6.5, 1.9, 'AVG CPC', '$2.99', 'May: $5.59')
tile(s, 8.5, 1.9, 'VIEWABILITY', '51.3%', 'May: 67.1%')
tile(s, 10.5, 1.9, 'CONVERSIONS', '11', 'view-through')
text(s, 0.5, 3.35, 6.0, 0.3, 'Month over Month', 11, True, NAVY)
table(s, 0.5, 3.65, 6.2, 2.2, [
    ['Metric', 'June', 'May', 'MoM'],
    ['Spend', '$33,784', '$26,890', '+26%'],
    ['Impressions', '41.96M', '16.48M', '+155%'],
    ['CPM', '$0.81', '$1.63', '-50%'],
    ['Viewability', '51.3%', '67.1%', '-15.8pp'],
    ['Conversions', '11 VT', '1', '+10'],
], col_widths=[2.0, 1.4, 1.4, 1.4])
card(s, 7.0, 3.35, 5.8, 3.55, 'READ', [
    'We scaled reach on purpose: 42M impressions and 18.2M devices '
    'reached as CPM fell by half to $0.81.',
    'The cost of cheap reach is quality. Viewability fell to 51%, below '
    'the 70% standard. We delivered a 49-site blocklist covering 32% of '
    'June spend on this line and recommend a campaign-level viewability '
    'floor.',
    'All 11 June conversions are view-through. Treat them as directional '
    'support, not proven response.',
], BLUE)

# ---- Slide 11: Q3 in Review — Spend ----------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'Q3 FY2026 IN REVIEW  ·  SPEND',
       'The program scaled every month of the quarter.', 3)
table(s, 0.5, 1.9, 6.6, 2.4, [
    ['Channel', 'Apr', 'May', 'Jun', 'Q3'],
    ['Bing Ads', '$15,289', '$15,972', '$25,659', '$56,920'],
    ['Meta', '$5,224', '$6,626', '$25,924', '$37,774'],
    ['Quantcast', '$15,005', '$26,890', '$33,784', '$75,678'],
    ['Azerion', '$10,777', '$29,302', '$35,026', '$75,105'],
    ['Total', '$46,294', '$78,790', '$120,393', '$245,477'],
], col_widths=[1.9, 1.15, 1.15, 1.2, 1.2], bold_last=True)
text(s, 0.5, 4.5, 6.6, 0.9, [
    'Month-over-month growth: +70% into May, +53% into June. June ran at '
    '2.6x the April level.',
], 9.5, False, BODY)
card(s, 7.5, 1.9, 5.3, 5.0, 'READ', [
    'Q3 closed at $245,477 in working media, approximately 24% of the '
    'FY2026 budget.',
    'The mix rebalanced as the program scaled. April leaned on search and '
    'display; June is a balanced four-channel program with no channel '
    'above 29% of spend.',
    'Every channel grew while holding or improving its efficiency where '
    'we can measure it, which is what earns the next stage of scale.',
], BLUE)

# ---- Slide 12: Q3 in Review — Applications ----------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'Q3 FY2026 IN REVIEW  ·  APPLICATIONS',
       'Application volume scaled with spend. Blended cost per '
       'application halved.', 4)
table(s, 0.5, 1.9, 6.6, 2.7, [
    ['Channel', 'Apr', 'May', 'Jun', 'Q3'],
    ['Bing Ads', '—', '33', '50', '83'],
    ['Azerion', '19', '37', '42', '98'],
    ['Quantcast', '0', '1', '11', '12'],
    ['Meta', '1', '0', 'Held', '1'],
    ['Blended applications', '20', '71', '103', '194'],
    ['Blended cost / application', '$2,315', '$1,110', '$1,169', '$1,265'],
], col_widths=[2.6, 0.95, 0.95, 1.0, 1.1], bold_last=True)
text(s, 0.5, 4.85, 6.6, 0.7, [
    'Search tracking went live in May; display lines count submitted '
    'applications, the reach line counts view-through, and paid-social '
    'conversions are under validation.',
], 7.5, False, SUB)
card(s, 7.5, 1.9, 5.3, 5.0, 'READ', [
    'June delivered 5x the application volume of April on 2.6x the '
    'spend.',
    'Blended cost per application fell from $2,315 in April to $1,169 in '
    'June and held near that level through a 53% spend increase, so the '
    'growth was bought at improving efficiency, not rising cost.',
    'Search and programmatic carried the measured volume (181 submitted '
    'applications). Paid social contributed reach and site traffic while '
    'its conversion validation completes.',
], GREEN)

# ---- Slide 13: Q3 Blended Outcomes vs Plan ----------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'Q3 FY2026 IN REVIEW  ·  BLENDED OUTCOMES VS PLAN',
       'The gap to plan sits in the organic base, not paid delivery.', 5)
table(s, 0.5, 1.9, 6.6, 2.7, [
    ['Metric (organic + paid)', 'Q3 Actual', 'Q3 Budget', 'vs Plan'],
    ['Live Apps Submitted', '2,167', '3,984', '-46%'],
    ['New Approved', '606', '957', '-37%'],
    ['New Funded', '106', '240', '-56%'],
    ['New Traded', '100', '180', '-44%'],
    ['Actives (monthly avg)', '399', '485', '-18%'],
    ['Trading Volume', '$4.3B', '$6.7B', '-36%'],
], col_widths=[2.8, 1.3, 1.3, 1.2])
text(s, 0.5, 4.85, 6.6, 0.5, [
    'Source: Budget vs Actual dashboard (non-cohort reforecast), GGMI, '
    'organic and paid blended. July excluded (full-month forecast).',
], 7.5, False, SUB)
card(s, 7.5, 1.9, 5.3, 5.0, 'READ', [
    'Blended submissions ran about 45% under plan and declined through '
    'the half-year (1,029 in January to 684 in June) while paid-driven '
    'applications grew every month of Q3. The shortfall sits in the '
    'non-paid base, which the next two slides trace to organic search.',
    'The funded stage shows the second constraint: -56% to plan, the '
    'application-to-funded step the onboarding roadmap is built to fix.',
    'Trading volume tracked the softer funnel at $4.3B against a $6.7B '
    'plan.',
], BLUE)

# ---- Slide 14: Site Traffic (GA4) ------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'SITE TRAFFIC  ·  SIX-MONTH VIEW (GA4, MEXICO)',
       'Traffic recovered in June. Media bought the recovery.', 6)
line_chart(s, 0.5, 1.95, 7.0, 2.85,
           [('Sessions', [9380, 22229, 12614, 7577, 5751, 9236]),
            ('Unique visitors', [4651, 14172, 6554, 3366, 2592, 5838])],
           [CHART_BLUE, CHART_GREEN], label_series=1)
card(s, 0.5, 5.0, 7.0, 1.9, 'THE HONEST READ', [
    'February was one paid flight, not organic growth; it ended in early '
    'March and is not a baseline. The March-May slide is organic search '
    "decline, confirmed in Search Console. June's rebound was bought: of "
    'the +3,485 sessions vs May, paid search added +2,601, paid social '
    '+537, display +274. Organic added +94.',
], BLUE)
card(s, 7.9, 1.95, 4.9, 4.95, 'WHAT IT MEANS', [
    'Media is doing its job. June proves it can move traffic on demand.',
    'Unique visitors are up 26% January to June, the legitimate bright '
    'spot in the half-year.',
    'The organic foundation needs its own workstream. Until it recovers, '
    'total site traffic will track media spend rather than compound on '
    'top of it.',
], GREEN)

# ---- Slide 15: Organic Search — the SEO issue --------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'ORGANIC SEARCH  ·  THE STRUCTURAL ISSUE',
       'Organic search fell 75% in six months. Media cannot buy it back.', 7)
text(s, 0.5, 1.9, 6.6, 0.25, 'Organic Search sessions — GA4, Mexico', 9,
     True, NAVY)
line_chart(s, 0.5, 2.2, 6.6, 2.3,
           [('Organic Search sessions', [4958, 3193, 3385, 2055, 1260,
                                         1313])],
           [CHART_BLUE], label_series=0)
table(s, 0.5, 4.75, 6.6, 1.7, [
    ['Search Console (Mexico)', 'January', 'June', 'Change'],
    ['Organic clicks', '4,811', '1,210', '-75%'],
    ['Demo account page (position)', '12', '23', 'slid'],
    ['Trading academy page (position)', '14', '24', 'slid'],
    ['Homepage, brand (position)', '—', '2.5', 'improved'],
], col_widths=[3.0, 1.2, 1.2, 1.2])
card(s, 7.5, 1.9, 5.3, 2.45, 'WHAT IS DRIVING IT', [
    'Dated news and analysis content lost 87-99% of its clicks as it '
    'aged out with nothing replacing it, and mid-funnel commercial pages '
    'slid in rank while brand improved. The decline spans every country '
    'the Spanish site serves: a site-level content issue, confirmed '
    'independently in site analytics.',
], RED)
card(s, 7.5, 4.55, 5.3, 2.45, 'WHAT WE RECOMMEND', [
    'Open a dedicated SEO workstream: restore the news and analysis '
    'publishing cadence, refresh the decayed assets, and run a '
    'mid-funnel rank-recovery plan (demo, academy, platform pages), '
    'measured on organic clicks and positions. Until it lands, paid '
    'media works harder each month just to hold total traffic flat.',
], GREEN)

# ---- Slide 16: Cross-Channel Priorities + Next Steps ------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'CROSS-CHANNEL PRIORITIES + NEXT STEPS',
       'Concentrate in-market. Finish the measurement. Start the SEO '
       'workstream.', 17)
card(s, 0.5, 1.9, 4.0, 4.4, 'PRIORITY', [
    'Complete the Mexico-only search restriction, with roughly $12.6K '
    'per month at stake.',
    'Run the paid-social placement audit and exclusions before '
    'reinstating its conversion reporting.',
], RED)
card(s, 4.8, 1.9, 4.0, 4.4, 'SCALING & OPPORTUNITY', [
    'Hold the budget mix pending the measurement fixes, then revisit '
    'allocation with clean July attribution.',
    'Search efficiency held at $513 and earns the first incremental '
    'dollar once spend is concentrated in Mexico.',
    'Open the SEO workstream for the Spanish site and apply the '
    'programmatic blocklist with a viewability floor.',
], GREEN)
card(s, 9.0, 1.9, 3.8, 4.4, 'MEASUREMENT TO CLOSE', [
    'Link the search account in GA4. 27% of June Mexico sessions are '
    'unattributed because of it, and July reporting will restate paid '
    'search larger once linked. We will flag the restatement.',
    'Complete the placement audit and pixel-event review, then reinstate '
    'paid-social conversions.',
], BLUE)
rect(s, 0.5, 6.5, 12.3, 0.55, GRAY_BG)
rect(s, 0.5, 6.5, 12.3, 0.08, BLUE)
text(s, 0.7, 6.63, 11.9, 0.35,
     'Concentrate spend on measurable, in-market delivery while the '
     'quality fixes land.', 10, True, NAVY)

# ---- Slide (created last, positioned 2): Summary — blended view -------------
s = prs.slides.add_slide(BLANK)
header(s, 'SUMMARY  ·  BLENDED VIEW (ORGANIC + PAID)',
       'Paid momentum built through Q3 while the blended base softened.', 2)
text(s, 0.5, 2.0, 3.1, 4.8, [
    '•  June submitted applications came in at 684, down 7% MoM; the H1 '
    'peak was February (1,036).',
    '•  New approved rose 7% to 209 and the approval rate recovered to '
    '31%, the best since February.',
    '•  New funded dipped 11% to 32 and the fund rate to 15%. The '
    'application-to-funding step is the funnel constraint.',
    '•  New traded closed at 29; trading volume held at $1.1B.',
    '•  Working media reached $120,393 in June (+53% MoM), the biggest '
    'month of 2026, while paid-driven applications grew every month of '
    'Q3.',
], 9.5, False, BODY)
table(s, 3.9, 2.0, 8.9, 3.3, [
    ['GGMI — blended', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'MoM'],
    ['Working Media Spend', '$1,717', '$20,578', '$42,424', '$46,294',
     '$78,790', '$120,393', '+53%'],
    ['Submitted Applications', '1,029', '1,036', '900', '751', '732',
     '684', '-7%'],
    ['Approved Clients', '288', '337', '261', '202', '195', '209', '+7%'],
    ['Approval Rate', '28%', '33%', '29%', '27%', '27%', '31%', '+4pp'],
    ['New Funded Clients', '66', '59', '63', '38', '36', '32', '-11%'],
    ['Fund Rate', '23%', '18%', '24%', '19%', '18%', '15%', '-3pp'],
    ['Cost per Funded Client', '$26', '$349', '$673', '$1,218', '$2,189',
     '$3,762', '+72%'],
    ['New Traded Clients', '56', '52', '57', '35', '36', '29', '-19%'],
    ['Cost per Traded Client', '$31', '$396', '$744', '$1,323', '$2,189',
     '$4,151', '+90%'],
    ['Trading Volume', '$5.1B', '$2.5B', '$2.3B', '$2.1B', '$1.1B',
     '$1.1B', 'flat'],
], col_widths=[2.15, 0.97, 0.97, 0.97, 0.97, 0.97, 0.97, 0.92])
text(s, 3.9, 5.5, 8.9, 0.5, [
    'Funnel metrics are blended organic + paid. Working media per the '
    'budget tracker. MoM = June vs May.',
], 7.5, False, SUB)

# Reorder: title, Summary, Q3 section, traffic/SEO, June detail, close.
NEW_ORDER = [0, 16, 10, 11, 12, 13, 14, 1, 2, 3, 4, 5, 6, 7, 8, 9, 15]
sld_lst = prs.slides._sldIdLst
ids = list(sld_lst)
for el in ids:
    sld_lst.remove(el)
for i in NEW_ORDER:
    sld_lst.append(ids[i])

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f'Saved {os.path.abspath(OUT)} ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)')
