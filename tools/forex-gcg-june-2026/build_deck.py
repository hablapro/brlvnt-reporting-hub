#!/usr/bin/env python3
"""Build the GCG June 2026 + Q3 FY2026 QBR deck (PPTX working draft).

Mirrors the GGMI June deck system (tools/forex-june-2026/build_deck.py):
QBR-first order, blended Summary lead slide, line charts, per-channel
Performance + Read pairs. GCG differences per approved narrative v1:
no SEO-crisis section; rank-limited search story; funnel-conversion
constraint; Azerion + Native carried at tracker spend pending vendor
detail. Client-facing spend = client budget tracker (Meta $30,711 per
Renzo 2026-07-17).
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import (XL_CHART_TYPE, XL_LEGEND_POSITION,
                             XL_MARKER_STYLE, XL_LABEL_POSITION)

NAVY = RGBColor(0x1A, 0x2A, 0x3A)
BLUE = RGBColor(0x3B, 0x59, 0x98)
GRAY_BG = RGBColor(0xF3, 0xF3, 0xF3)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
BODY = RGBColor(0x32, 0x37, 0x3E)
MUTED = RGBColor(0x78, 0x82, 0x8C)
SUB = RGBColor(0x5A, 0x5A, 0x5A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
GRID = RGBColor(0xD9, 0xD9, 0xD9)

HERE = os.path.dirname(os.path.abspath(__file__))
LOGO = os.path.join(HERE, '..', 'forex-june-2026', 'assets',
                    'logo_1_fd6b9b.png')
OUT = os.path.join(HERE, '..', '..', 'report-client-decks',
                   '06. GCG_US_June_2026_Performance_Review.pptx')
BREADCRUMB = 'FOREX.com  |  GCG (US Hispanic)  ·  June 2026'
FOOTER_LEFT = 'Performance Review'
N_SLIDES = 15
MONTHS = ['Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun']

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def rect(slide, x, y, w, h, fill):
    sh = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, size=9.5, bold=False, color=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    paras = [runs] if isinstance(runs, str) else runs
    for i, ptext in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        if i > 0:
            p.space_before = Pt(6)
        r = p.add_run()
        r.text = ptext
        r.font.name = 'Arial'
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
    return tb


def header(slide, section, headline, pageno):
    rect(slide, 0, 0, 13.333, 0.5, NAVY)
    text(slide, 0.5, 0.1, 9.0, 0.3, BREADCRUMB, 9, True, WHITE)
    if os.path.exists(LOGO):
        slide.shapes.add_picture(LOGO, Inches(12.2), Inches(0.11),
                                 Inches(0.98), Inches(0.28))
    text(slide, 0.5, 0.62, 12.0, 0.3, section, 11, True, BLUE)
    text(slide, 0.5, 0.95, 12.3, 0.75, headline, 20, True, NAVY)
    rect(slide, 0, 7.2, 13.333, 0.3, GRAY_BG)
    text(slide, 0.5, 7.23, 6.0, 0.22, FOOTER_LEFT, 8, False, MUTED)
    text(slide, 12.1, 7.23, 0.9, 0.22, f'{pageno} / {N_SLIDES}', 8, False,
         MUTED, align=PP_ALIGN.RIGHT)


def tile(slide, x, y, label, value, sub=None, w=1.9):
    h = 1.1 if sub else 1.0
    rect(slide, x, y, w, h, GRAY_BG)
    rect(slide, x, y, 0.1, h, BLUE)
    text(slide, x + 0.2, y + 0.05, w - 0.3, 0.3, label, 8, True, MUTED)
    text(slide, x + 0.2, y + 0.32, w - 0.3, 0.5, value, 15, True, NAVY)
    if sub:
        text(slide, x + 0.2, y + 0.82, w - 0.3, 0.24, sub, 8, False, SUB)


def card(slide, x, y, w, h, title, body_paras, accent=BLUE):
    rect(slide, x, y, w, h, GRAY_BG)
    rect(slide, x, y, w, 0.1, accent)
    text(slide, x + 0.2, y + 0.2, w - 0.4, 0.3, title, 10, True, accent)
    text(slide, x + 0.2, y + 0.55, w - 0.4, h - 0.75, body_paras, 9.5,
         False, BODY)


def table(slide, x, y, w, h, rows, col_widths=None, bold_last=False):
    n_r, n_c = len(rows), len(rows[0])
    gfx = slide.shapes.add_table(n_r, n_c, Inches(x), Inches(y),
                                 Inches(w), Inches(h))
    tbl = gfx.table
    tbl.first_row = False
    tbl.horz_banding = False
    if col_widths:
        for ci, cw in enumerate(col_widths):
            tbl.columns[ci].width = Inches(cw)
    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = tbl.cell(ri, ci)
            cell.margin_left = cell.margin_right = Inches(0.06)
            cell.margin_top = cell.margin_bottom = Inches(0.03)
            cell.fill.solid()
            if ri == 0:
                cell.fill.fore_color.rgb = BLUE
            elif bold_last and ri == n_r - 1:
                cell.fill.fore_color.rgb = GRAY_BG
            else:
                cell.fill.fore_color.rgb = WHITE if ri % 2 else GRAY_BG
            p = cell.text_frame.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.name = 'Arial'
            r.font.size = Pt(8.5)
            r.font.bold = (ri == 0) or (ci == 0) or (bold_last and ri == n_r - 1)
            r.font.color.rgb = WHITE if ri == 0 else BODY
            if ci > 0:
                p.alignment = PP_ALIGN.RIGHT
    return gfx


def line_chart(slide, x, y, w, h, series, colors, label_series=None):
    cd = CategoryChartData()
    cd.categories = MONTHS
    for name, vals in series:
        cd.add_series(name, vals)
    gfx = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Inches(x),
                                 Inches(y), Inches(w), Inches(h), cd)
    ch = gfx.chart
    ch.has_title = False
    ch.font.size = Pt(8)
    ch.font.name = 'Arial'
    ch.font.color.rgb = BODY
    ch.has_legend = len(series) > 1
    if ch.has_legend:
        ch.legend.position = XL_LEGEND_POSITION.BOTTOM
        ch.legend.include_in_layout = False
        ch.legend.font.size = Pt(8)
    val_ax = ch.value_axis
    val_ax.has_major_gridlines = True
    val_ax.major_gridlines.format.line.color.rgb = GRID
    val_ax.major_gridlines.format.line.width = Pt(0.5)
    val_ax.tick_labels.number_format = '#,##0'
    val_ax.tick_labels.number_format_is_linked = False
    val_ax.format.line.fill.background()
    cat_ax = ch.category_axis
    cat_ax.has_major_gridlines = False
    cat_ax.format.line.color.rgb = GRID
    for i, s in enumerate(ch.series):
        s.smooth = False
        s.format.line.color.rgb = colors[i]
        s.format.line.width = Pt(2.25)
        s.marker.style = XL_MARKER_STYLE.CIRCLE
        s.marker.size = 6
        s.marker.format.fill.solid()
        s.marker.format.fill.fore_color.rgb = colors[i]
        s.marker.format.line.color.rgb = WHITE
        if label_series is not None and i == label_series:
            s.data_labels.show_value = True
            s.data_labels.number_format = '#,##0'
            s.data_labels.number_format_is_linked = False
            s.data_labels.font.size = Pt(7.5)
            s.data_labels.font.color.rgb = BODY
            s.data_labels.position = XL_LABEL_POSITION.ABOVE
    return gfx


# ---- S1: Title --------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, 13.333, 7.5, WHITE)
rect(s, 0, 6.4, 13.333, 1.1, NAVY)
rect(s, 0.5, 2.2, 0.1, 2.5, BLUE)
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(0.5), Inches(0.5), Inches(2.5),
                         Inches(0.71))
text(s, 0.9, 2.2, 11.0, 0.4, 'QUARTERLY BUSINESS REVIEW', 14, True, BLUE)
text(s, 0.9, 2.7, 12.0, 1.0, 'GCG (US Hispanic) — Paid Media', 36, True, NAVY)
text(s, 0.9, 3.8, 12.0, 0.6,
     'Q3 FY2026 Business Review  ·  June 2026 Performance', 22, False, NAVY)
text(s, 0.9, 4.5, 12.0, 0.4,
     'Google Ads · Meta · Azerion · Quantcast · Native', 12, False, SUB)
text(s, 0.9, 6.65, 12.0, 0.4,
     'Reporting period: June 1–30, 2026   |   Currency: USD   |   '
     'Prepared by Berelvant · 07.17.2026', 10, False, WHITE)

# ---- S2: Summary — blended view ---------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'SUMMARY  ·  BLENDED VIEW (ORGANIC + PAID)',
       'Volume is arriving. The start-to-submit step is the constraint.', 2)
text(s, 0.5, 2.0, 3.1, 4.8, [
    '•  June live applications rose 6% to 306 and submissions 7% to 322; '
    'the approval rate held at 45%.',
    '•  New funded held at 41 and new traded closed at 30. The fund step '
    '(28%) remains the gate to account growth.',
    '•  Working media reached $117,024 in June (+53% MoM), the biggest '
    'GCG month of 2026, with the Native pilot adding a fifth channel.',
    '•  Sessions grew 30% MoM on the paid scale-up while application '
    'starts stayed flat, so conversion, not volume, is the constraint.',
    '•  Q2 closed 14% below Q1 on applications (1,033 vs 1,203) while '
    'media scaled 5x quarter over quarter.',
], 9.5, False, BODY)
table(s, 3.9, 2.0, 8.9, 3.5, [
    ['GCG — blended', 'Jan', 'Feb', 'Mar', 'Apr', 'May', 'Jun', 'MoM'],
    ['Working Media Spend', '—', '$2,237', '$49,106', '$58,324', '$76,274',
     '$117,024', '+53%'],
    ['Unique Sessions', '4,187', '5,075', '31,785', '43,192', '52,137',
     '67,542', '+30%'],
    ['App Starts', '2,306', '1,871', '2,820', '2,767', '2,617', '2,590',
     '-1%'],
    ['Live Apps Submitted', '430', '344', '391', '398', '289', '306',
     '+6%'],
    ['Approved Clients', '202', '150', '172', '186', '134', '145', '+8%'],
    ['Approval Rate', '46%', '43%', '42%', '45%', '45%', '45%', '0pp'],
    ['New Funded Clients', '58', '46', '44', '65', '40', '41', '+3%'],
    ['Fund Rate', '29%', '31%', '26%', '35%', '30%', '28%', '-2pp'],
    ['Cost per Funded Client', '—', '$49', '$1,116', '$897', '$1,907',
     '$2,854', '+50%'],
    ['New Traded Clients', '45', '40', '33', '53', '31', '30', '-3%'],
    ['Cost per Traded Client', '—', '$56', '$1,488', '$1,100', '$2,460',
     '$3,901', '+59%'],
], col_widths=[2.15, 0.97, 0.97, 0.97, 0.97, 0.97, 0.97, 0.92])
text(s, 3.9, 5.7, 8.9, 0.5, [
    'Funnel metrics are blended organic + paid (client dashboard). '
    'Working media per the budget tracker. MoM = June vs May.',
], 7.5, False, SUB)

# ---- S3: Q3 Spend -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'Q3 FY2026 IN REVIEW  ·  SPEND',
       'The program doubled from April and added a fifth channel.', 3)
table(s, 0.5, 1.9, 6.6, 2.7, [
    ['Channel', 'Apr', 'May', 'Jun', 'Q3'],
    ['Google Ads', '$15,120', '$15,201', '$22,524', '$52,845'],
    ['Meta', '$12,167', '$12,243', '$30,711', '$55,121'],
    ['Quantcast', '$15,015', '$22,359', '$30,559', '$67,933'],
    ['Azerion', '$16,021', '$26,472', '$29,586', '$72,079'],
    ['Native', '—', '—', '$3,645', '$3,645'],
    ['Total', '$58,324', '$76,274', '$117,024', '$251,621'],
], col_widths=[1.9, 1.15, 1.15, 1.2, 1.2], bold_last=True)
text(s, 0.5, 4.8, 6.6, 0.9, [
    'Month-over-month growth: +31% into May, +53% into June. June ran at '
    '2x the April level.',
], 9.5, False, BODY)
card(s, 7.5, 1.9, 5.3, 5.0, 'READ', [
    'Q3 closed at $251,621 in working media, roughly 16% of the FY2026 '
    'budget; the program is at 19.35% of budget year to date.',
    'The mix stayed balanced as it scaled: no channel above 27% of June '
    'spend, with search, social, and two programmatic lines all growing '
    'and the Native pilot launching as the fifth channel.',
    'Paid social led June growth (+151%), followed by search (+48%) and '
    'the reach line (+37%).',
], BLUE)

# ---- S4: Q3 Applications ----------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'Q3 FY2026 IN REVIEW  ·  APPLICATIONS',
       'Search and display carried the measured volume.', 4)
table(s, 0.5, 1.9, 6.6, 2.7, [
    ['Channel', 'Apr', 'May', 'Jun', 'Q3'],
    ['Google Ads', '12', '76', '67', '155'],
    ['Azerion', '76', '43', '58', '177'],
    ['Quantcast', '2', '8', '15', '25'],
    ['Meta (pixel events)*', '73', '109', '136', '318'],
    ['Measured submitted apps', '88', '119', '125', '332'],
], col_widths=[2.6, 0.95, 0.95, 1.0, 1.1], bold_last=True)
text(s, 0.5, 4.85, 6.6, 0.9, [
    '* Meta counts pixel application events, mostly application starts, '
    'on a traffic objective; not comparable to submitted applications. '
    'Quantcast is view-through, directional only. Measured submitted '
    'apps = Google + Azerion.',
], 7.5, False, SUB)
card(s, 7.5, 1.9, 5.3, 5.0, 'READ', [
    'The quarter delivered 332 measured submitted applications, with '
    'June its best month (125).',
    'Azerion returned to form in June: 58 applications, up 35% on a 12% '
    'budget increase, at a $510 CPA (down 17%).',
    'June exposed the search ceiling: spend rose 48% and applications '
    'fell 12%, because the account is rank-limited (64-76% of available '
    'impressions lost to ad rank, only 9-13% to budget). July search '
    'work is bids, quality, and ad strength before budget.',
], GREEN)

# ---- S5: Q1 vs Q2 blended funnel ---------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'Q1 VS Q2  ·  BLENDED FUNNEL (ORGANIC + PAID)',
       'Four times the traffic, the same funded accounts.', 5)
table(s, 0.5, 1.9, 6.6, 3.1, [
    ['Metric', 'Q1 (Jan-Mar)', 'Q2 (Apr-Jun)', 'Change'],
    ['Unique Sessions', '41,047', '162,871', '+297%'],
    ['App Starts', '6,997', '7,974', '+14%'],
    ['Submitted', '1,203', '1,033', '-14%'],
    ['Live Apps Submitted', '1,165', '993', '-15%'],
    ['Approved', '524', '465', '-11%'],
    ['New Funded', '148', '146', '-1%'],
    ['New Traded', '118', '114', '-3%'],
], col_widths=[2.4, 1.4, 1.4, 1.2])
text(s, 0.5, 5.25, 6.6, 0.5, [
    'Client dashboard, calendar quarters, blended organic + paid.',
], 7.5, False, SUB)
card(s, 7.5, 1.9, 5.3, 5.0, 'READ', [
    'Traffic quadrupled quarter over quarter while application starts '
    'grew 14% and submissions fell 14%. Each funnel stage dampened the '
    'volume the media bought: the start rate diluted from 17% to 5% and '
    'the submit rate from 17% to 13%.',
    'Funded and traded held flat, so the business outcome absorbed none '
    'of the scale yet.',
    'The work that converts this traffic is funnel work: the '
    'start-to-submit step first, then the approved-to-funded step. '
    'Media is filling the top; the next accounts come from conversion.',
], BLUE)

# ---- S6: Site traffic chart ---------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'SITE TRAFFIC  ·  SPANISH-LANGUAGE AUDIENCE (GA4, US)',
       'A stable base with a paid June lift. No organic problem here.', 6)
line_chart(s, 0.5, 1.95, 7.0, 2.85,
           [('Sessions', [69895, 58481, 74971, 69553, 69927, 80231]),
            ('Unique visitors', [24928, 20841, 29355, 30251, 29808, 41795])],
           [BLUE, GREEN], label_series=1)
card(s, 0.5, 5.0, 7.0, 1.9, 'THE READ', [
    'Spanish-language sessions on the US property held near 70K from '
    'January through May and rose 15% in June with the media scale-up; '
    'unique visitors rose 40% to 41,795. The traffic base under GCG is '
    'healthy and growing.',
], BLUE)
card(s, 7.9, 1.95, 4.9, 4.95, 'WHAT IT MEANS', [
    'Media is expanding the audience: June brought the highest '
    'Spanish-language visitor count of the year.',
    'Measurement is working: site analytics capture roughly 57% of paid '
    'social clicks as sessions, and search reconciles exactly.',
    'The gap between traffic growth and application growth is a '
    'conversion problem, not a traffic or measurement problem, which is '
    'why the funnel work on the previous slide leads the '
    'recommendations.',
], GREEN)

# ---- S7: Exec summary (June) --------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'JUNE 2026  ·  EXECUTIVE SUMMARY',
       'The biggest GCG media month of 2026, with one channel at its '
       'ceiling.', 7)
table(s, 0.5, 1.9, 7.4, 2.3, [
    ['Channel', 'Spend', 'Impr', 'Clicks', 'Conv', 'CPA'],
    ['Google Ads', '$22,524', '181,470', '12,888', '67', '$336'],
    ['Meta', '$30,711', '3.06M', '74,572', '136*', '—'],
    ['Quantcast', '$30,559', '37.2M', '2,942', '15 VT*', '—'],
    ['Azerion', '$29,586', '4.52M', '19,653', '58', '$510'],
    ['Native', '$3,645', 'pilot', '—', '—', '—'],
    ['Total', '$117,024', '—', '—', '—', '—'],
], col_widths=[2.33, 1.06, 1.06, 1.06, 0.85, 1.06], bold_last=True)
text(s, 0.5, 4.35, 7.4, 0.45,
     '* Meta counts pixel application events (mostly starts) on a traffic '
     'objective. Programmatic view-through is directional only. '
     'Conversions use different events per channel and are never summed.',
     7.5, False, SUB)
card(s, 8.1, 1.9, 4.7, 2.4, 'HEADLINES', [
    'June was the biggest GCG media month of 2026: $117,024 across five '
    'lines, up 53% from May. Azerion led conversion with 58 applications '
    'at $510 (+35%), search added 67 at $336, and the Native pilot '
    'entered the market.',
], GREEN)
card(s, 0.5, 4.95, 7.4, 2.0, 'WHAT THE DATA SAYS', [
    'Search hit its rank ceiling: +48% spend bought -12% applications. '
    'Paid social scaled 151% with healthy delivery (71% of clicks became '
    'landing-page views). The reach line scaled 37% and its viewability '
    'fell to 46.9%, which the blocklist addresses.',
], BLUE)
card(s, 8.1, 4.95, 4.7, 2.0, 'JUNE PRIORITY', [
    'Shift paid social to the conversion objective, run the search '
    'ad-rank program before adding budget, and apply the programmatic '
    'blocklist with a viewability floor.',
], RED)

# ---- S8: Cont' exec -----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, "CONT' EXECUTIVE SUMMARY",
       'Three moves shape July.', 8)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT IS STILL FORMING', [
    'Search efficiency at scale. June proved the account is rank-limited: '
    'every campaign loses 64-76% of available impressions to ad rank and '
    'only 9-13% to budget. More budget at current ad rank buys worse '
    'auctions.',
    'Paid social conversion measurement. Delivery and site capture are '
    'healthy, but the campaign still runs a traffic objective, so its '
    '136 pixel events are mostly application starts rather than '
    'submitted applications.',
    'The Native pilot. $3,645 entered the market in June; the first '
    'delivery and cost-per-application read lands with vendor '
    'reporting.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Complete the paid-social shift to a conversion objective and judge '
    'the channel on submitted applications from July.',
    'Run the search ad-rank program on the proven TrackB terms: bid and '
    'quality work plus an RSA refresh, then scale into the recovered '
    'impression share.',
    'Apply the 18-site programmatic blocklist (32% of June spend on that '
    'line) and set a campaign-level viewability floor.',
    'Put the start-to-submit funnel step on the joint roadmap with the '
    'client; media is filling the top of the funnel.',
], GREEN)

# ---- S9: Google perf ----------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'GOOGLE ADS  ·  PERFORMANCE',
       '67 submitted applications at $336. The account hit its rank '
       'ceiling.', 9)
tile(s, 0.5, 1.9, 'SPEND', '$22,524', '+48% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '181,470', '4 campaigns')
tile(s, 4.5, 1.9, 'CLICKS', '12,888', 'CTR 7.10%')
tile(s, 6.5, 1.9, 'AVG CPC', '$1.75', 'May: $1.53')
tile(s, 8.5, 1.9, 'SUBMITTED APPS', '67', 'May: 76')
tile(s, 10.5, 1.9, 'CPA', '$336', 'May: $200')
text(s, 0.5, 3.35, 7.0, 0.3, 'Campaigns', 11, True, NAVY)
table(s, 0.5, 3.65, 7.6, 1.7, [
    ['Campaign', 'Spend', 'Clicks', 'Apps', 'CPA', 'Search IS'],
    ['TrackA Trust', '$7,796', '5,712', '18', '$433', '27%'],
    ['TrackB Authority', '$6,826', '3,198', '23', '$297', '24%'],
    ['TrackB Platform', '$4,489', '2,568', '17', '$264', '13%'],
    ['Brand Search', '$3,413', '1,410', '9', '$379', '11%'],
], col_widths=[2.4, 1.1, 1.1, 0.8, 1.0, 1.2])
card(s, 8.4, 3.35, 4.4, 3.55, 'HIGHLIGHTS', [
    'All 67 conversions are submitted applications (Step 5 goal, '
    'definition unchanged from May).',
    'TrackB stayed the efficiency core: Platform $264 and Authority '
    '$297. TrackA Trust converted at $433.',
    'Spend rose 48% but applications fell 12%, because extra budget at '
    'current ad rank bought worse auctions, not more volume.',
], GREEN)
text(s, 0.5, 5.55, 7.6, 1.3, [
    'Impression share tells the story: campaigns lose 64-76% of '
    'available impressions to ad rank and only 9-13% to budget. Brand '
    'holds just 11% of its own auctions. The account has a large runway '
    'that budget alone cannot buy.',
], 9.5, False, BODY)

# ---- S10: Google read/rec ------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'GOOGLE ADS  ·  READ & RECOMMENDATIONS',
       'Earn rank first. Then scale.', 10)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT THE DATA SHOWS', [
    'May was the efficiency peak: 76 applications at $200 with even '
    'weekly pacing. June added 48% more budget and got 12% fewer '
    'applications at $336.',
    'The auction data isolates the cause: lost impression share is '
    'rank-driven everywhere (64-76%), not budget-driven (9-13%). The '
    'account cannot buy its way past ad rank at current bids, quality, '
    'and ad strength.',
    'TrackB remains the proven converter set ($264-$297); Brand at 11% '
    'impression share is under-defended.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Run the ad-rank program before adding budget: bid up the proven '
    'TrackB terms, refresh RSAs to lift ad strength, and rebuild quality '
    'score on the Trust track.',
    'Defend Brand: 11% impression share on brand terms leaves the '
    'highest-intent auctions to competitors.',
    'Hold June budget levels while rank recovers, then scale into the '
    'impression share the rank work opens.',
], GREEN)

# ---- S11: Meta perf ------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'META  ·  PERFORMANCE',
       'Scaled 2.5x with clean delivery. Ready for the conversion '
       'objective.', 11)
tile(s, 0.5, 1.9, 'SPEND', '$30,711', '+151% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '3.06M', 'CPM $10.04')
tile(s, 4.5, 1.9, 'REACH', '1.53M', 'freq 2.0')
tile(s, 6.5, 1.9, 'LINK CLICKS', '74,572', 'CTR 2.44%')
tile(s, 8.5, 1.9, 'LP VIEWS', '52,639', '70.6% of clicks')
tile(s, 10.5, 1.9, 'PIXEL EVENTS', '136', 'starts rollup')
text(s, 0.5, 3.35, 6.0, 0.3, 'Ad Sets', 11, True, NAVY)
table(s, 0.5, 3.65, 6.4, 1.4, [
    ['Ad set', 'Spend', 'Link clicks', 'LPV', 'Pixel events'],
    ['trackA_pros_us_es', '$16,924', '35,710', '25,078', '84'],
    ['trackB_pros_us_es', '$17,787', '38,862', '27,561', '52'],
], col_widths=[2.2, 1.1, 1.2, 1.0, 1.2])
card(s, 7.2, 3.35, 5.6, 3.55, 'HIGHLIGHTS', [
    'Delivery held through a 2.5x scale-up: landing-page views ran at '
    '70.6% of link clicks (May: 65.9%) and site analytics capture '
    'roughly 57% of clicks as sessions. The measurement chain works.',
    'The campaign still runs the Q2 traffic objective, so its 136 pixel '
    'events are mostly application starts. The conversion-objective '
    'shift is the unlock for judging this channel on submitted '
    'applications.',
], GREEN)
text(s, 0.5, 5.35, 6.4, 1.3, [
    'June audience delivery reached 1.53M people at 2.0 frequency, the '
    'largest Spanish-language reach of any GCG channel with verified '
    'site arrival.',
], 9.5, False, BODY)

# ---- S12: Meta read/rec ---------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'META  ·  READ & RECOMMENDATIONS',
       'Make the objective match the goal.', 12)
card(s, 0.5, 1.9, 6.1, 5.1, 'WHAT THE DATA SHOWS', [
    'The channel does what a traffic objective asks: cheap, verified '
    'site visitors at scale. What it does not do yet is optimize toward '
    'applications, because the objective never asked it to.',
    'Every delivery signal supports the shift: stable CPM through the '
    'scale-up, improving click-to-LPV rate, healthy analytics capture, '
    'and a growing retargeting pool from four months of traffic.',
], BLUE)
card(s, 6.9, 1.9, 5.9, 5.1, 'WHAT WE RECOMMEND', [
    'Complete the shift to a conversion objective in July and judge the '
    'channel on submitted applications.',
    'Seed the conversion campaigns with the pixel and the four-month '
    'retargeting pool before scaling prospecting.',
    'Keep the trackA/trackB split so creative learnings carry over from '
    'the search program.',
], GREEN)

# ---- S13: Quantcast ------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'QUANTCAST  ·  PERFORMANCE & READ',
       'A deliberate reach month, with the same quality bill to manage.',
       13)
tile(s, 0.5, 1.9, 'SPEND', '$30,559', '+37% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '37.2M', '+38% MoM')
tile(s, 4.5, 1.9, 'CPM', '$0.82', 'May: $0.83')
tile(s, 6.5, 1.9, 'DEVICE REACH', '23.2M', 'freq 1.6')
tile(s, 8.5, 1.9, 'VIEWABILITY', '46.9%', '70% standard')
tile(s, 10.5, 1.9, 'RESULTS', '15', '13 view-through')
card(s, 0.5, 3.35, 6.1, 3.55, 'READ', [
    'The reach line scaled on purpose: 37.2M impressions to 23.2M '
    'devices at a flat $0.82 CPM.',
    'Viewability fell to 46.9%, well under the 70% standard. We '
    'delivered an 18-site blocklist covering $9,752 (32%) of June spend '
    'on this line and recommend a campaign-level viewability floor.',
    'Results rose 8 to 15, but 13 of 15 are view-through: directional '
    'support, not proven response.',
], BLUE)
card(s, 6.9, 3.35, 5.9, 3.55, 'TOP BLOCKLIST ENTRIES', [
    'yahoo.com $2,270 at 15.7% viewability; clarin.com $1,315 at 26%; '
    'weather.com $1,278 at 17.7%; genius.com $771 at 0.9%; zillow.com '
    '$685 at 5.6%.',
    'The list ships with this report; applying it plus a viewability '
    'floor moves the line toward the standard without giving up the '
    'reach efficiency.',
], RED)

# ---- S14: Azerion (+ Native line) --------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'AZERION  ·  PERFORMANCE & READ',
       '58 applications, up 35% on a 12% budget increase. Fewer starts, '
       'better starts.', 14)
tile(s, 0.5, 1.9, 'SPEND', '$29,586', '+12% MoM')
tile(s, 2.5, 1.9, 'IMPRESSIONS', '4.52M', 'CPM $6.55')
tile(s, 4.5, 1.9, 'CLICKS', '19,653', 'CTR 0.44%')
tile(s, 6.5, 1.9, 'APP STARTS', '447', 'May: 1,398')
tile(s, 8.5, 1.9, 'APPLICATIONS', '58', '$510 CPA')
tile(s, 10.5, 1.9, 'VIEWABILITY', '58.8%', '70% standard')
text(s, 0.5, 3.35, 6.0, 0.3, 'Ad Sets (Applications)', 11, True, NAVY)
table(s, 0.5, 3.65, 6.0, 2.3, [
    ['Ad Set', 'Spend', 'Apps', 'CPA'],
    ['Trusted Broker', '$4,933', '14', '$352'],
    ['Broker 1', '$4,964', '13', '$382'],
    ['Language Broker', '$4,900', '10', '$490'],
    ['Professional Tools', '$4,940', '9', '$549'],
    ['Spanish Platform', '$4,933', '7', '$705'],
    ['Trust HTML', '$4,916', '5', '$983'],
], col_widths=[2.6, 1.1, 0.9, 1.2])
card(s, 6.8, 3.35, 6.0, 3.0, 'READ', [
    'Applications rose 35% (43 to 58) at a $510 CPA, down 17%: the '
    'strongest converter of the June program.',
    'Application starts fell 1,398 to 447 while completion improved '
    'from 3.1% to 13.0%, which is the start-to-submit fix May called '
    'for: fewer, better-qualified starts.',
    'Delivery is confirmed US-only (state-level geo report). '
    'Viewability at 58.8% sits under the 70% standard; the 300x600 '
    'format converts best ($124 vendor-basis) and 728x90 worst.',
], GREEN)
rect(s, 6.8, 6.5, 6.0, 0.55, GRAY_BG)
rect(s, 6.8, 6.5, 6.0, 0.08, BLUE)
text(s, 7.0, 6.62, 5.6, 0.35,
     'Native pilot: $3,645 in month one; first delivery read lands with '
     'vendor reporting.', 9, False, BODY)

# ---- S15: Close -----------------------------------------------------------------
s = prs.slides.add_slide(BLANK)
header(s, 'CROSS-CHANNEL PRIORITIES + NEXT STEPS',
       'Convert the traffic. Earn the rank. Guard the quality.', 15)
card(s, 0.5, 1.9, 4.0, 4.4, 'PRIORITY', [
    'Shift paid social to the conversion objective and judge it on '
    'submitted applications.',
    'Put the start-to-submit funnel step on the joint roadmap; media is '
    'filling the top of the funnel and conversion is the constraint.',
], RED)
card(s, 4.8, 1.9, 4.0, 4.4, 'SCALING & OPPORTUNITY', [
    'Run the search ad-rank program (bids, quality, RSA refresh) on '
    'TrackB terms, defend Brand, then scale into the recovered '
    'impression share.',
    'Hold the reach line at pace once the blocklist and viewability '
    'floor land.',
], GREEN)
card(s, 9.0, 1.9, 3.8, 4.4, 'TO CLOSE THE QUARTER', [
    'Carry the Azerion start-quality fix forward: completion tripled to '
    '13% in June; concentrate budget on Trusted Broker and Broker 1 and '
    'push viewability over the standard.',
    'First Native pilot read lands with vendor reporting; judged on '
    'cost per submitted application.',
], BLUE)
rect(s, 0.5, 6.5, 12.3, 0.55, GRAY_BG)
rect(s, 0.5, 6.5, 12.3, 0.08, BLUE)
text(s, 0.7, 6.63, 11.9, 0.35,
     'Media is filling the funnel. July converts it.', 10, True, NAVY)

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print(f'Saved {os.path.abspath(OUT)} ({len(prs.slides._sldIdLst)} slides)')
