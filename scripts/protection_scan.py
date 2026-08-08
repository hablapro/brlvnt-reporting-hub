#!/usr/bin/env python3
"""Protection scan — the programmatic half of the pre-delivery gate.

docs/DOCTRINE.md requires this scan before any client deliverable leaves the
shop, and says the forbidden-vocabulary check must be "verified
programmatically." It was run by eye through the June cycle and self-indicting
phrasing survived a manual stop-slop pass anyway. This is that check.

    python3 scripts/protection_scan.py <file> [<file> ...]
    python3 scripts/protection_scan.py report-client-decks/*.pptx

Reads .pptx, .docx, .xlsx, .md and .txt. Exits 1 if any BLOCK finding lands,
0 otherwise. WARN findings need a human read, not an automatic stop.

Scope: CLIENT-FACING artifacts only, meaning the deck, the performance Sheet
and any doc that leaves the shop. Internal files legitimately use every
forbidden word, so running this on docs/, qa/ or a build status produces a
wall of BLOCKs that mean nothing. Point it at deliverables.

What it cannot do: judge whether a number matches the client's tracker (that
is scripts/verify_numbers.py), spot a competitor name it has never seen, or
tell you a slide is missing. A clean run is a floor, not an approval.
"""
import re
import sys
import zipfile
from pathlib import Path

# ---- BLOCK: internal vocabulary that must never reach a client artifact ----
# From docs/DOCTRINE.md. These words expose bookkeeping, contract or
# compliance mechanics the client materials deliberately do not discuss.
FORBIDDEN = {
    'fee': 'billing mechanics — client materials show fee-inclusive figures',
    'raw': 'implies an adjusted figure exists',
    'adjustment': 'reconciliation detail is internal-only',
    'adjusted': 'reconciliation detail is internal-only',
    'reconciled': 'reconciliation detail is internal-only',
    'reconciliation': 'reconciliation detail is internal-only',
    'basis': 'signals a restated basis the client did not ask for',
    'violation': 'never characterise our own delivery as a violation',
    'breach': 'never characterise our own delivery as a breach',
    'compliance': 'compliance-adjacent framing needs a human decision',
    'internal': 'internal/external boundary must not be visible',
    'cheap': 'say "low CPM" or give the figure (docs/DOCTRINE.md voice rules)',
    'wasted': 'frames our own spend as waste',
    'blended cpa': 'never publish a blended CPA across different events',
}

# ---- BLOCK: statements against interest ------------------------------------
# Sentences where the grammatical actor buying or causing a bad outcome is us.
# Stop-slop catches AI tells; it does not catch this. All four seed patterns
# are real phrases that shipped in a June draft.
AGAINST_INTEREST = [
    (r'\b(bought|buying|purchased)\s+\w*\s*(worse|poor|low[- ]quality|bad)',
     'we are the actor buying a bad outcome'),
    (r'\b(our|we|the campaigns?)\s+\w{0,12}\s*(bought|skewed|drove down|'
     r'degraded|hurt)\b',
     'we are the actor causing a negative result'),
    (r'\b(extra|additional|more)\s+budget\s+bought\b',
     'frames added spend as buying a worse result'),
    (r'\bfor most of (h1|h2|the year|the quarter)\b',
     'time-scope creep: do not extend a bad fact backward'),
    (r'\bwe (failed|missed|did not|didn\'t)\b',
     'agency accountability belongs in internal notes, not the deck'),
    (r'\b(we|berelvant) (should have|could have)\b',
     'agency accountability belongs in internal notes, not the deck'),
]

# ---- WARN: needs a human read ---------------------------------------------
# Third element restricts a rule to certain extensions. Every pattern here was
# tightened against the shipped June GGMI deck until it stopped crying wolf:
# a bare \bonly\b fired 11 times on "Mexico-only" targeting, and a noisy gate
# is a gate people learn to skip.
WARN = [
    (r'\b(cost per (funded|traded|approved)|funded (client|account)|'
     r'deposits?)\b',
     'downstream metric we do not control: add context or replace with cost '
     'per submitted application', None),
    (r'\b(pending|TBC|TBD|awaiting|not yet (known|available))\b',
     'reads as a hole; restate as a plan with an owner and a date', None),
    # superlatives only when they are actually superlative claims
    (r'\b(the only|the first|the best|most trusted|el único|el primero|'
     r'el mejor|más confiable)\b',
     'absolute or superlative claim — GCG compliance guardrails apply', None),
    (r'\b(safe|protected|guaranteed?|protegido|garantiza)\b',
     'safety or guarantee language — use "regulated" / "regulado"', None),
    (r'\b(build wealth|construir riqueza|diversify your income)\b',
     'income or wealth promise — GCG compliance guardrails forbid it', None),
    # the em-dash rule is a prose rule. The decks use an em dash as the
    # established title separator ("GGMI (LATAM) — Paid Media"), so scanning
    # .pptx for it produces 20 hits and zero decisions.
    ('—', 'em dash — voice rules use commas, periods or a line break',
     {'.md', '.txt', '.docx'}),
]


def text_from(path: Path):
    """Yield (locator, text) pairs. Locator is a slide/sheet/line reference."""
    ext = path.suffix.lower()
    if ext == '.pptx':
        from pptx import Presentation
        for i, slide in enumerate(Presentation(str(path)).slides, 1):
            for sh in slide.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip():
                    yield f'slide {i}', sh.text_frame.text
                if getattr(sh, 'has_table', False) and sh.has_table:
                    for row in sh.table.rows:
                        for c in row.cells:
                            if c.text.strip():
                                yield f'slide {i} (table)', c.text
    elif ext == '.docx':
        # read the XML directly so python-docx is not a dependency
        with zipfile.ZipFile(path) as z:
            xml = z.read('word/document.xml').decode('utf8', 'replace')
        xml = re.sub(r'</w:p>', '\n', xml)
        for i, line in enumerate(re.sub(r'<[^>]+>', '', xml).split('\n'), 1):
            if line.strip():
                yield f'para {i}', line
    elif ext == '.xlsx':
        import openpyxl
        wb = openpyxl.load_workbook(str(path), data_only=True)
        for ws in wb.worksheets:
            for row in ws.iter_rows():
                for cell in row:
                    if isinstance(cell.value, str) and cell.value.strip():
                        yield f'{ws.title}!{cell.coordinate}', cell.value
    elif ext in ('.md', '.txt'):
        for i, line in enumerate(
                path.read_text(encoding='utf8', errors='replace').split('\n'),
                1):
            if line.strip():
                yield f'line {i}', line
    else:
        raise ValueError(f'unsupported file type: {path.suffix}')


def scan(path: Path):
    findings = []
    ext = path.suffix.lower()
    for loc, chunk in text_from(path):
        low = chunk.lower()
        for word, why in FORBIDDEN.items():
            # word-boundary match so "basis" does not fire inside "basispoint"
            # and "raw" does not fire inside "drawn"
            if re.search(rf'\b{re.escape(word)}\b', low):
                findings.append(('BLOCK', loc, word, why, chunk))
        for pat, why in AGAINST_INTEREST:
            m = re.search(pat, low)
            if m:
                findings.append(('BLOCK', loc, m.group(0), why, chunk))
        for pat, why, exts in WARN:
            if exts and ext not in exts:
                continue
            m = re.search(pat, chunk, re.I)
            if m:
                findings.append(('WARN', loc, m.group(0).strip(), why, chunk))
    return findings


def main(argv):
    if not argv:
        print(__doc__)
        return 2
    blocked = 0
    for arg in argv:
        path = Path(arg)
        if not path.exists():
            print(f'!! not found: {path}')
            blocked += 1
            continue
        try:
            findings = scan(path)
        except Exception as exc:                       # noqa: BLE001
            print(f'!! {path.name}: {type(exc).__name__}: {exc}')
            blocked += 1
            continue

        blocks = [f for f in findings if f[0] == 'BLOCK']
        warns = [f for f in findings if f[0] == 'WARN']
        print(f'\n=== {path.name} — {len(blocks)} BLOCK, {len(warns)} WARN')
        for level, loc, hit, why, chunk in blocks + warns:
            excerpt = ' '.join(chunk.split())
            if len(excerpt) > 150:
                excerpt = excerpt[:147] + '...'
            print(f'  [{level}] {loc}: "{hit}" — {why}')
            print(f'          {excerpt}')
        blocked += len(blocks)

    print(f'\n{"FAIL" if blocked else "PASS"}: {blocked} blocking finding(s).')
    if blocked:
        print('A BLOCK is a stop, not a suggestion. Fix the copy or, if the '
              'word is genuinely required, get it ruled on and add the '
              'exception to docs/DOCTRINE.md — not to this script silently.')
    return 1 if blocked else 0


if __name__ == '__main__':
    sys.exit(main(sys.argv[1:]))
